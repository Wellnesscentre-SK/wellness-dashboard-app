"""
Normal Yearly Report — merged annual analysis (e.g. "2026 Data Analysis").

Accepts EITHER a single period dict (backward compatible) OR a list of
monthly period dicts. When a list is given, all months (Jan–Dec) are merged
into one annual dataset: overall totals, averages, percentages, month-over-
month trends and key changes are computed automatically and rendered with
one clear visualization per insight — no duplicate charts.

Design system, alignment, colors and typography match the other reports.
"""

import io
import re
from datetime import date

from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import config as C
import components as CM
import template_slides as TS


# ═══════════════════════════════════════════════════════════════════════════════
# MERGING + ANALYSIS
# ═══════════════════════════════════════════════════════════════════════════════

def _as_list(periods) -> list:
    if isinstance(periods, dict):
        return [periods]
    return [p for p in periods if isinstance(p, dict)]


def _month_sort_key(p):
    start = p.get("start")
    if start:
        try:
            y, m, d = str(start)[:10].split("-")
            return (int(y), int(m), int(d))
        except (ValueError, TypeError):
            pass
    return (9999, 12, 31)


def _short_month(p):
    start = p.get("start")
    if start:
        try:
            return date.fromisoformat(str(start)[:10]).strftime("%b")
        except ValueError:
            pass
    label = str(p.get("label", ""))
    m = re.match(r"([A-Za-z]{3})\w*", label)
    return m.group(1) if m else label[:6] or "?"


def _detect_year(months):
    years = set()
    for p in months:
        start = p.get("start")
        if start:
            years.add(str(start)[:4])
    if len(years) == 1:
        return years.pop()
    for p in months:
        m = re.search(r"(20\d{2})", str(p.get("label", "")))
        if m:
            return m.group(1)
    return str(date.today().year)


def _merge(months: list) -> dict:
    agg = {
        "new": 0, "followup": 0, "grand": 0,
        "gender": {}, "mode": {}, "referral": {},
        "concern": {}, "stakeholder": {},
        "vertical": {v: {"new": 0, "followup": 0, "total": 0} for v in C.VERT_KEYS},
        "by_vertical": {},
    }
    for p in months:
        agg["new"] += int(p.get("new", 0) or 0)
        agg["followup"] += int(p.get("followup", 0) or 0)
        agg["grand"] += int(p.get("grand", 0) or 0)
        for dim in ("gender", "mode", "referral", "concern", "stakeholder"):
            for k, v in (p.get(dim) or {}).items():
                agg[dim][k] = agg[dim].get(k, 0) + int(v or 0)
        combined = C.combine_verticals(p.get("vertical") or {})
        for vk in C.VERT_KEYS:
            pv = combined.get(vk) or {}
            for sub in ("new", "followup", "total"):
                agg["vertical"][vk][sub] += int(pv.get(sub, 0) or 0)
        for src, dims in (p.get("by_vertical") or {}).items():
            tgt = agg["by_vertical"].setdefault(src, {})
            for dim, cats in dims.items():
                t = tgt.setdefault(dim, {})
                for k, v in cats.items():
                    t[k] = t.get(k, 0) + int(v or 0)
    return agg


def _pct(part, whole):
    return (part / whole * 100) if whole else 0.0


def _change_pct(old, new):
    if not old:
        return None
    return (new - old) / old * 100


def _generate_insights(merged, months):
    """Auto-computed annual performance summary lines."""
    grand = merged["grand"]
    n_months = max(1, len(months))
    lines = []

    lines.append(
        f"Annual performance: {grand:,} total cases across {len(months)} "
        f"month{'s' if len(months) != 1 else ''} (average {grand // n_months:,} cases/month)."
    )
    lines.append(
        f"Case split: New {merged['new']:,} ({_pct(merged['new'], grand):.0f}%) | "
        f"Follow-up {merged['followup']:,} ({_pct(merged['followup'], grand):.0f}%)."
    )

    if len(months) >= 2:
        totals = [(m.get("label") or _short_month(m), int(m.get("grand", 0) or 0))
                  for m in months]
        peak_label, peak_val = max(totals, key=lambda t: t[1])
        low_label, low_val = min(totals, key=lambda t: t[1])
        lines.append(f"Peak month: {peak_label} ({peak_val:,} cases); "
                     f"lowest: {low_label} ({low_val:,} cases).")

        biggest = None
        for (la, va_), (lb, vb_) in zip(totals, totals[1:]):
            delta = vb_ - va_
            pct = _change_pct(va_, vb_)
            if biggest is None or abs(delta) > abs(biggest[2]):
                biggest = (la, lb, delta, pct)
        if biggest and biggest[2]:
            la, lb, delta, pct = biggest
            word = "jump" if delta > 0 else "drop"
            pct_txt = f" ({pct:+.1f}%)" if pct is not None else ""
            lines.append(f"Key change: biggest month-over-month {word} "
                         f"{la} to {lb} ({delta:+,}{pct_txt}).")

        half = len(totals) // 2
        if half:
            h1 = sum(v for _, v in totals[:half])
            h2 = sum(v for _, v in totals[half:])
            pct = _change_pct(h1, h2)
            if pct is not None:
                direction = "up" if h2 >= h1 else "down"
                lines.append(f"Trend: second half is {direction} "
                             f"{abs(pct):.1f}% vs first half ({h1:,} to {h2:,}).")

    vk = C.VERT_KEYS
    vert_totals = [(C.VERTICALS[v], merged["vertical"][v]["total"]) for v in vk]
    top_vert, top_vert_val = max(vert_totals, key=lambda t: t[1])
    lines.append(f"Top vertical: {top_vert} with {top_vert_val:,} cases "
                 f"({_pct(top_vert_val, grand):.0f}% of annual load).")

    def top(dim, labels):
        vals = [(lbl, merged[dim].get(lbl, 0)) for lbl in labels]
        lbl, val = max(vals, key=lambda t: t[1]) if vals else ("—", 0)
        return lbl, val

    c_lbl, c_val = top("concern", C.CONCERN_LABELS)
    lines.append(f"Top concern addressed: {c_lbl} ({c_val:,} cases, "
                 f"{_pct(c_val, grand):.0f}% of all cases).")
    s_lbl, s_val = top("stakeholder", C.STAKEHOLDER_LABELS)
    lines.append(f"Largest stakeholder group: {s_lbl} ({s_val:,} cases, "
                 f"{_pct(s_val, grand):.0f}%).")
    return lines


# ═══════════════════════════════════════════════════════════════════════════════
# DECK BUILD
# ═══════════════════════════════════════════════════════════════════════════════

def build(periods, insights: list = None) -> bytes:
    months = sorted(_as_list(periods), key=_month_sort_key)
    if not months:
        raise ValueError("normal_yearly.build requires at least one period")
    merged = _merge(months)
    year = _detect_year(months)
    label = f"{year} Data Analysis"
    grand = merged["grand"]
    new_n = merged["new"]
    fu_n = merged["followup"]
    avg_month = round(grand / max(1, len(months)))
    p = merged
    vk = C.VERT_KEYS

    prs = CM.create_presentation()

    # ── SLIDE 1: COVER ────────────────────────────────────────────────────────
    CM.add_cover(prs, "Annual Wellness Data Report", label)

    # ── SLIDE 2: KPI CARDS + VERTICALS TABLE + SINGLE PIE ─────────────────────
    s2 = CM.add_slide(prs)
    CM.add_logo(s2)
    CM.add_kpi_cards(s2, [
        ("Total Cases", grand, C.TOTAL_CLR),
        ("New Cases", new_n, C.NEW_CASE_CLR),
        ("Follow-up Cases", fu_n, C.FOLLOWUP_CLR),
        (f"Avg / Month ({len(months)} mo)", avg_month, C.TOP_VERT_CLR),
    ])

    v_labels = [C.VERTICALS[v] for v in vk]
    v_totals = [p["vertical"][v]["total"] for v in vk]
    v_new = [p["vertical"][v]["new"] for v in vk]
    v_fu = [p["vertical"][v]["followup"] for v in vk]

    vtable = [["Vertical", "Total Cases", "% Share"]]
    for i, v in enumerate(vk):
        vtable.append([C.VERTICALS[v], str(v_totals[i]),
                        f"{_pct(v_totals[i], grand):.0f}%"])
    CM.add_table(s2, Inches(0.25), Inches(1.85), Inches(5.5), Inches(1.5), vtable)
    CM.add_pie_chart(s2, Inches(6.5), Inches(1.70), Inches(6.5), Inches(4.5),
                     v_labels, v_totals)

    # ── SLIDE 3: MONTHLY TREND (ONE CHART, JAN–DEC) ───────────────────────────
    s3 = CM.add_slide(prs)
    CM.add_logo(s3)
    CM.add_section_header(s3, f"MONTHLY TREND  |  {label}", Inches(0.15), Inches(0.75),
                          Inches(13.0), Inches(0.45), bg_color=C.PRIMARY)
    trend_labels = [_short_month(m) for m in months]
    trend_values = [int(m.get("grand", 0) or 0) for m in months]
    CM.add_column_chart(s3, Inches(0.15), Inches(1.35), Inches(13.0), Inches(5.5),
                        trend_labels, [("Total Cases", trend_values)],
                        legend_position="bottom",
                        series_colors=[C.TOTAL_CLR], show_percentage=False)

    # ── SLIDE 4: NEW vs FOLLOW-UP BY VERTICAL ─────────────────────────────────
    s4 = CM.add_slide(prs)
    CM.add_logo(s4)
    CM.add_section_header(s4, "NEW vs FOLLOW-UP BY VERTICAL", Inches(0.15), Inches(0.75),
                          Inches(13.0), Inches(0.45), bg_color=C.PRIMARY)
    nf_table = [["Vertical", "New Cases", "Follow-up", "New %", "F/up %"]]
    for i, v in enumerate(vk):
        t = v_totals[i] or 1
        nf_table.append([C.VERTICALS[v], str(v_new[i]), str(v_fu[i]),
                         f"{v_new[i]/t*100:.0f}%", f"{v_fu[i]/t*100:.0f}%"])
    CM.add_table(s4, Inches(0.25), Inches(1.35), Inches(6.0), Inches(2.0), nf_table)
    CM.add_column_chart(s4, Inches(6.8), Inches(1.20), Inches(6.2), Inches(5.5),
                        v_labels, [("New", v_new), ("Follow-up", v_fu)],
                        legend_position="bottom",
                        series_colors=[C.NEW_CASE_CLR, C.FOLLOWUP_CLR])

    # ── SLIDE 5: DEMOGRAPHICS & MODE TABLES ───────────────────────────────────
    s5 = CM.add_slide(prs)
    CM.add_logo(s5)
    CM.add_section_header(s5, "DEMOGRAPHICS & MODE OF SESSION", Inches(0.15), Inches(0.75),
                          Inches(13.0), Inches(0.45), bg_color=C.PRIMARY)

    g_table = [["Gender", "Count", "% Share"]]
    for g in C.GENDER_LABELS:
        gval = p["gender"].get(g, 0)
        g_table.append([g, str(gval), f"{_pct(gval, grand):.0f}%"])

    m_table = [["Mode", "Count", "% Share"]]
    for m in C.MODE_LABELS:
        mval = p["mode"].get(m, 0)
        m_table.append([m, str(mval), f"{_pct(mval, grand):.0f}%"])

    CM.add_table(s5, Inches(0.25), Inches(1.35), Inches(4.0), Inches(2.2), g_table)
    CM.add_table(s5, Inches(4.50), Inches(1.35), Inches(4.0), Inches(2.2), m_table)

    r_table = [["Referral", "Count", "% Share"]]
    for r in C.REFERRAL_LABELS:
        rval = p["referral"].get(r, 0)
        r_table.append([r, str(rval), f"{_pct(rval, grand):.0f}%"])
    CM.add_table(s5, Inches(8.75), Inches(1.35), Inches(4.3), Inches(3.0), r_table)

    # ── SLIDE 6: CONCERN TABLE + SINGLE PIE ───────────────────────────────────
    s6 = CM.add_slide(prs)
    CM.add_logo(s6)
    CM.add_section_header(s6, "RANGE OF CONCERN ADDRESSED", Inches(0.15), Inches(0.75),
                          Inches(13.0), Inches(0.45), bg_color=C.PRIMARY)
    c_labels = C.CONCERN_LABELS
    c_values = [p["concern"].get(c, 0) for c in c_labels]
    c_table = [["Concern", "Count", "% Share"]]
    for i, c in enumerate(c_labels):
        c_table.append([c, str(c_values[i]), f"{_pct(c_values[i], grand):.0f}%"])
    CM.add_table(s6, Inches(0.25), Inches(1.35), Inches(6.0), Inches(4.5), c_table)
    CM.add_pie_chart(s6, Inches(6.8), Inches(1.35), Inches(6.2), Inches(5.0),
                     c_labels, c_values)

    # ── SLIDE 7: STAKEHOLDER TABLE + SINGLE PIE ───────────────────────────────
    s7 = CM.add_slide(prs)
    CM.add_logo(s7)
    CM.add_section_header(s7, "STAKEHOLDER", Inches(0.15), Inches(0.75),
                          Inches(13.0), Inches(0.45), bg_color=C.PRIMARY)
    st_labels = C.STAKEHOLDER_LABELS
    st_values = [p["stakeholder"].get(s, 0) for s in st_labels]
    st_table = [["Stakeholder", "Count", "% Share"]]
    for i, s in enumerate(st_labels):
        st_table.append([s, str(st_values[i]), f"{_pct(st_values[i], grand):.0f}%"])
    CM.add_table(s7, Inches(0.25), Inches(1.35), Inches(6.0), Inches(4.5), st_table)
    CM.add_pie_chart(s7, Inches(6.8), Inches(1.35), Inches(6.2), Inches(5.0),
                     st_labels, st_values)

    # ── SLIDE 8: GENDER DISTRIBUTION (SINGLE PIE) ─────────────────────────────
    s8 = CM.add_slide(prs)
    CM.add_logo(s8)
    CM.add_section_header(s8, f"GENDER DISTRIBUTION  |  {label}", Inches(0.15), Inches(0.75),
                          Inches(13.0), Inches(0.45), bg_color=C.PRIMARY)
    g_labels = C.GENDER_LABELS
    g_values = [p["gender"].get(g, 0) for g in g_labels]
    CM.add_pie_chart(s8, Inches(0.3), Inches(1.35), Inches(12.7), Inches(5.5),
                     g_labels, g_values)

    # ── SLIDE 9: MODE OF SESSION (SINGLE PIE) ─────────────────────────────────
    s9 = CM.add_slide(prs)
    CM.add_logo(s9)
    CM.add_section_header(s9, f"MODE OF SESSION  |  {label}", Inches(0.15), Inches(0.75),
                          Inches(13.0), Inches(0.45), bg_color=C.PRIMARY)
    m_labels = C.MODE_LABELS
    m_values = [p["mode"].get(m, 0) for m in m_labels]
    CM.add_pie_chart(s9, Inches(0.3), Inches(1.35), Inches(12.7), Inches(5.5),
                     m_labels, m_values)

    # ── SLIDES 10-12: TEMPLATE (GENDER / CONCERN / STAKEHOLDER, vertical-wise) ─
    TS.add_template_slides(prs, p)

    # ── SLIDE 13: YEARLY KEY INSIGHTS (AUTO-COMPUTED SUMMARY) ─────────────────
    s13 = CM.add_slide(prs)
    CM.add_logo(s13)
    CM.add_section_header(s13, f"YEARLY KEY INSIGHTS  |  {label}", Inches(0.15), Inches(0.75),
                          Inches(13.0), Inches(0.45), bg_color=C.PRIMARY)
    auto_insights = _generate_insights(merged, months)
    y = Inches(1.40)
    for i, line in enumerate(auto_insights[:9], 1):
        CM.add_text(s13, f"{i}.", Inches(0.6), y, Inches(0.4), Inches(0.55),
                    font_size=C.BODY_FONT, bold=True, color=C.PRIMARY,
                    alignment=PP_ALIGN.RIGHT)
        CM.add_text(s13, line, Inches(1.1), y, Inches(11.9), Inches(0.55),
                    font_size=Pt(12), bold=False, color=C.BLACK)
        y += Inches(0.58)

    # ── SLIDE 14: AI INSIGHTS ─────────────────────────────────────────────────
    if insights is None:
        import ai_analyzer
        insights = ai_analyzer.analyze(p)
    CM.add_insights_slide(prs, "AI DATA INSIGHTS", insights, label)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
