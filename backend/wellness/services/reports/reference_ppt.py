"""Reference-style PPT modules used by every comparative PPT export.

Converts Django Period/CaseRow models to the dict format expected by the
standalone ppt_generator engine, then delegates to that engine for PPT
generation. This keeps the visual output identical between standalone
and webapp usage.

Design source: "Report to asc dean_29_JULY_2026 - WO.pptx"
"""

from __future__ import annotations

import io
import os
import sys
from datetime import date

# ═══════════════════════════════════════════════════════════════════════════════
# IMPORT STANDALONE PPT ENGINE
# The ppt_generator lives at the project root (sibling of backend/).
# ═══════════════════════════════════════════════════════════════════════════════

_PPT_ROOT = os.path.normpath(os.path.join(
    os.path.dirname(__file__), "..", "..", "..", "..", "ppt_generator"
))
if _PPT_ROOT not in sys.path:
    sys.path.insert(0, _PPT_ROOT)

from wellness.services.reports import ppt as legacy  # noqa: E402 — data helpers
from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402


def _import_ppt_module(mod_name):
    """Import a ppt_generator module without Django's config namespace collision.

    Evicts ALL ppt_generator modules (config, components, weekly, monthly …)
    from sys.modules so that they always re-import against ppt_generator/config.py
    instead of Django's config module.  The Django config keys are restored
    afterward to leave the rest of the server unaffected.
    """
    import importlib

    # Modules that live inside ppt_generator and use `import config as C`
    _PPT_MODULES = {
        "config", "components", "weekly", "monthly", "yearly",
        "normal_week", "normal_monthly", "normal_yearly",
        "template_slides", "runner",
    }

    # 1. Save + evict Django's config so it can't bleed into ppt_generator imports.
    saved_django_cfg = {}
    for key in list(sys.modules):
        if key == "config" or key.startswith("config."):
            saved_django_cfg[key] = sys.modules.pop(key)

    # 2. Evict any previously cached ppt_generator modules so they are re-imported
    #    fresh (this forces `components.py` to re-execute its module-level
    #    `import config as C`, which now picks up ppt_generator/config.py).
    evicted_ppt = {}
    for key in list(sys.modules):
        if key in _PPT_MODULES:
            evicted_ppt[key] = sys.modules.pop(key)

    try:
        mod = importlib.import_module(mod_name)
        return mod
    finally:
        # 3. Restore Django's config (never the ppt_generator ones – let them
        #    be re-evicted on the next call so the fix is always consistent).
        sys.modules.update(saved_django_cfg)


# ═══════════════════════════════════════════════════════════════════════════════
# CONSTANTS (mirror the standalone config for positions)
# ═══════════════════════════════════════════════════════════════════════════════

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
RED = RGBColor(0xFF, 0x00, 0x00)
BLACK = RGBColor(0x00, 0x00, 0x00)
GREY = RGBColor(0x80, 0x80, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
VERTICALS = ("WC", "TA", "YD", "MW")
VERTICAL_LABELS = ("WC", "Team A", "YD", "MW")
VERTICAL_COLORS = (
    RGBColor(0x44, 0x72, 0xC4),
    RGBColor(0xED, 0x7D, 0x31),
    RGBColor(0xA5, 0xA5, 0xA5),
    RGBColor(0xFF, 0xC0, 0x00),
)
MARGIN = Inches(0.4)
HEADER_TOP = Inches(0.16)
HEADER_H = Inches(0.46)
CHART_TOP = Inches(1.35)
CHART_H = Inches(5.55)
LEFT = Inches(0.35)
RIGHT = Inches(6.72)
CHART_W = Inches(6.25)
FOOTER_TOP = Inches(7.15)

# ═══════════════════════════════════════════════════════════════════════════════
# PERIOD → DICT CONVERTER
# ═══════════════════════════════════════════════════════════════════════════════

def _period_to_dict(period) -> dict:
    """Convert a Django Period model instance to the ppt_generator dict format."""
    rows = legacy.index_rows(period)

    def total(ct, field):
        return sum(int(getattr(rows[(ct, v)], field, 0) or 0) for v in VERTICALS)

    def cat_total(fields):
        return sum(
            int(getattr(rows[(ct, v)], f, 0) or 0)
            for ct in ("new", "followup")
            for v in VERTICALS
            for f in fields
        )

    new_n = total("new", "total_cases")
    fu_n = total("followup", "total_cases")

    gender_fields = ["gender_male", "gender_female", "gender_other"]
    mode_fields = ["mode_online", "mode_in_person", "mode_phone"]
    referral_fields = ["referral_self", "referral_director", "referral_dean",
                       "referral_friend", "referral_mitr"]
    concern_fields = [
        "concern_anxiety", "concern_stress", "concern_career", "concern_interpersonal",
        "concern_self_dev", "concern_clinical", "concern_addiction", "concern_medical",
        "concern_suicidal",
    ]
    stake_fields = [
        "stake_ug", "stake_pg", "stake_phd", "stake_dual", "stake_faculty",
        "stake_employee_family", "stake_postdoc", "stake_unidentified",
    ]

    gender_vals = [total("new", f) + total("followup", f) for f in gender_fields]
    mode_vals = [total("new", f) + total("followup", f) for f in mode_fields]
    referral_vals = [total("new", f) + total("followup", f) for f in referral_fields]
    concern_vals = [total("new", f) + total("followup", f) for f in concern_fields]
    stake_vals = [total("new", f) + total("followup", f) for f in stake_fields]

    vert_dict = {}
    by_vertical = {}
    for v, label in zip(VERTICALS, VERTICAL_LABELS):
        rn = rows[("new", v)]
        rf = rows[("followup", v)]

        def _sum(f, rn=rn, rf=rf):
            return int(getattr(rn, f, 0) or 0) + int(getattr(rf, f, 0) or 0)

        v_new = int(getattr(rn, "total_cases", 0) or 0)
        v_fu = int(getattr(rf, "total_cases", 0) or 0)
        vert_dict[v] = {"new": v_new, "followup": v_fu, "total": v_new + v_fu}
        by_vertical[v] = {
            "gender": dict(zip(legacy.GENDERS, [_sum(f) for f in gender_fields])),
            "mode": dict(zip(legacy.MODES, [_sum(f) for f in mode_fields])),
            "referral": dict(zip(legacy.REFERRALS, [_sum(f) for f in referral_fields])),
            "concern": dict(zip(legacy.CONCERNS, [_sum(f) for f in concern_fields])),
            "stakeholder": dict(zip(legacy.STAKEHOLDERS, [_sum(f) for f in stake_fields])),
        }

    label = f"{_date_label(period.period_start)} to {_date_label(period.period_end)}"

    return {
        "label": label,
        "start": str(period.period_start),
        "end": str(period.period_end),
        "new": new_n,
        "followup": fu_n,
        "grand": new_n + fu_n,
        "gender": dict(zip(legacy.GENDERS, gender_vals)),
        "mode": dict(zip(legacy.MODES, mode_vals)),
        "referral": dict(zip(legacy.REFERRALS, referral_vals)),
        "concern": dict(zip(legacy.CONCERNS, concern_vals)),
        "stakeholder": dict(zip(legacy.STAKEHOLDERS, stake_vals)),
        "vertical": vert_dict,
        "by_vertical": by_vertical,
    }


def _date_label(d: date) -> str:
    suffix = "th" if 10 <= d.day % 100 <= 20 else {1: "st", 2: "nd", 3: "rd"}.get(d.day % 10, "th")
    return f"{d.day}{suffix} {d.strftime('%b %Y')}"


# ═══════════════════════════════════════════════════════════════════════════════
# BACKWARD-COMPAT API  (same function signatures as before)
# ═══════════════════════════════════════════════════════════════════════════════

def _ensure_dict(x):
    """Accept a Period instance OR an already-converted dict."""
    return x if isinstance(x, dict) else _period_to_dict(x)


def _ensure_dicts(xs):
    return [_ensure_dict(p) for p in xs]


def build_weekly_comparison(period_a, period_b, source_label="uploaded Wellness Excel reports") -> bytes:
    """Build 10-slide weekly comparison using the new engine."""
    try:
        mod = _import_ppt_module("weekly")
        build_weekly = mod.build
        a = _ensure_dict(period_a)
        b = _ensure_dict(period_b)
        return build_weekly(a, b)
    except Exception:
        a = _ensure_dict(period_a)
        b = _ensure_dict(period_b)
        return _build_weekly_fallback(a, b, source_label)


def build_monthly_comparison(period_a, period_b, insights=None,
                             source_label="uploaded Wellness Excel reports") -> bytes:
    """Build monthly comparison using the new engine."""
    try:
        mod = _import_ppt_module("monthly")
        build_monthly = mod.build
        a = _ensure_dict(period_a)
        b = _ensure_dict(period_b)
        key_insights = None
        if insights:
            key_insights = [item["text"] if isinstance(item, dict) else str(item)
                            for item in (insights if isinstance(insights, list) else [])]
        return build_monthly(a, b, key_insights=key_insights)
    except Exception:
        a = _ensure_dict(period_a)
        b = _ensure_dict(period_b)
        return _build_monthly_fallback(a, b, insights, source_label)


def build_yearly(periods_a, periods_b,
                 fy1_label="FY 2024-25", fy2_label="FY 2025-26") -> bytes:
    """Build 12-slide yearly comparison using the new engine."""
    try:
        mod = _import_ppt_module("yearly")
        build_yearly_fn = mod.build
        a_list = _ensure_dicts(periods_a)
        b_list = _ensure_dicts(periods_b)
        return build_yearly_fn(a_list, b_list, fy1_label, fy2_label)
    except Exception:
        a_list = _ensure_dicts(periods_a)
        b_list = _ensure_dicts(periods_b)
        return _build_yearly_fallback(a_list, b_list, fy1_label, fy2_label)


def build_normal_week(period, insights=None) -> bytes:
    """Build single-week wellness report using the new engine."""
    try:
        mod = _import_ppt_module("normal_week")
        build_nw = mod.build
        data = _ensure_dict(period)
        return build_nw(data, insights=insights)
    except Exception:
        data = _ensure_dict(period)
        return _build_normal_week_fallback(data)


def build_normal_monthly(period, insights=None) -> bytes:
    """Build single-month wellness report using the new engine."""
    try:
        mod = _import_ppt_module("normal_monthly")
        build_nm = mod.build
        data = _ensure_dict(period)
        return build_nm(data, insights=insights)
    except Exception:
        data = _ensure_dict(period)
        return _build_normal_monthly_fallback(data)


def build_normal_yearly(period, insights=None) -> bytes:
    """Merged annual report: every monthly period of the calendar year (Jan-Dec)
    is combined into one annual analysis. Accepts a single Period or a list."""
    try:
        mod = _import_ppt_module("normal_yearly")
        build_ny = mod.build
        if isinstance(period, (list, tuple)):
            pool = _ensure_dicts(period)
        else:
            from wellness.models import Period as PeriodModel

            year = period.period_start.year
            pool = _ensure_dicts(
                PeriodModel.objects.filter(
                    report_type=PeriodModel.ReportType.MONTHLY,
                    period_start__year=year,
                    superseded_by__isnull=True,
                ).order_by("period_start")
            )
            if not pool:
                pool = [_ensure_dict(period)]
        return build_ny(pool, insights=insights)
    except Exception:
        if isinstance(period, (list, tuple)):
            pool = _ensure_dicts(period)
        else:
            pool = [_ensure_dict(period)]
        return _build_normal_yearly_fallback(pool)


# ═══════════════════════════════════════════════════════════════════════════════
# FALLBACK BUILDERS  (used if ppt_generator is not on sys.path)
# ═══════════════════════════════════════════════════════════════════════════════

def _text(slide, left, top, width, height, value, size=12, *, color=BLACK,
          bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(value)
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _new_slide(prs, title="Wellness Data - Report"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    _text(slide, MARGIN, HEADER_TOP, Inches(5.8), HEADER_H, title, 24,
          color=RED, bold=True)
    return slide


def _footer(slide, page, source="Data source: uploaded Wellness Excel report"):
    _text(slide, MARGIN, FOOTER_TOP, Inches(10.4), Inches(0.2), source, 8, color=GREY)
    _text(slide, Inches(12.2), FOOTER_TOP, Inches(0.7), Inches(0.2), f"Page {page}", 8,
          color=GREY, align=PP_ALIGN.RIGHT)


def _section(slide, title, left_label, right_label=None):
    _text(slide, MARGIN, Inches(0.72), Inches(12.5), Inches(0.35), title, 16,
          color=BLACK, bold=True, align=PP_ALIGN.CENTER)
    _text(slide, LEFT, Inches(1.08), CHART_W, Inches(0.25), left_label, 11,
          color=RED, bold=True, align=PP_ALIGN.CENTER)
    if right_label is not None:
        _text(slide, RIGHT, Inches(1.08), CHART_W, Inches(0.25), right_label, 11,
              color=RED, bold=True, align=PP_ALIGN.CENTER)


def _chart(slide, chart_type, left, top, width, height, labels, series, title="", colors=None):
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_LEGEND_POSITION
    data = CategoryChartData()
    data.categories = [str(label) for label in labels]
    for name, values in series:
        data.add_series(name, [int(value or 0) for value in values])
    frame = slide.shapes.add_chart(chart_type, left, top, width, height, data)
    chart = frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(8)
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_category_name = chart_type.name == "PIE"
    dl.show_value = True
    dl.show_percentage = False
    dl.number_format = "0"
    dl.font.size = Pt(8)
    if chart_type.name == "PIE":
        palette = colors or VERTICAL_COLORS
        for i, point in enumerate(chart.series[0].points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = palette[i % len(palette)]
    else:
        for i, s in enumerate(chart.series):
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = (colors or VERTICAL_COLORS)[i % len(colors or VERTICAL_COLORS)]
    return frame


def _table(slide, left, top, width, height, rows):
    from pptx.enum.text import PP_ALIGN
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), left, top, width, height)
    table = table_shape.table
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RED if ri == 0 else WHITE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                p.font.name = "Calibri"
                p.font.size = Pt(9)
                p.font.bold = ri == 0
                p.font.color.rgb = WHITE if ri == 0 else BLACK
    return table_shape


def _pair_pies(prs, page, title, left_label, right_label, labels, left_values, right_values, source):
    from pptx.enum.chart import XL_CHART_TYPE
    slide = _new_slide(prs)
    _section(slide, title, left_label, right_label)
    _chart(slide, XL_CHART_TYPE.PIE, LEFT, CHART_TOP, CHART_W, CHART_H,
           labels, [(title, left_values)], colors=VERTICAL_COLORS)
    _chart(slide, XL_CHART_TYPE.PIE, RIGHT, CHART_TOP, CHART_W, CHART_H,
           labels, [(title, right_values)], colors=VERTICAL_COLORS)
    _footer(slide, page, source)
    return slide


def _insights(a, b, labels):
    def change(label, before, after):
        delta = after - before
        return f"{label}: {before} to {after} ({delta:+d})."
    moves = [
        change("Grand total cases", a["grand"], b["grand"]),
        change("New cases", a["new"], b["new"]),
        change("Follow-up cases", a["followup"], b["followup"]),
    ]
    largest = max(range(len(labels)), key=lambda i: b["vertical_total"][i] - a["vertical_total"][i])
    smallest = min(range(len(labels)), key=lambda i: b["vertical_total"][i] - a["vertical_total"][i])
    moves.append(change(f"Vertical with largest increase ({labels[largest]})", a["vertical_total"][largest], b["vertical_total"][largest]))
    moves.append(change(f"Vertical with largest change ({labels[smallest]})", a["vertical_total"][smallest], b["vertical_total"][smallest]))
    for key, name, category_labels in (
        ("gender", "Gender", legacy.GENDERS),
        ("mode", "Mode", legacy.MODES),
    ):
        index = max(range(len(category_labels)), key=lambda i: abs(b[key][i] - a[key][i]))
        moves.append(change(f"Largest {name} movement ({category_labels[index]})", a[key][index], b[key][index]))
    return moves


def _build_weekly_fallback(period_a, period_b, source_label):
    """Fallback weekly builder using dict data directly."""
    a = _ensure_dict(period_a) if not isinstance(period_a, dict) else period_a
    b = _ensure_dict(period_b) if not isinstance(period_b, dict) else period_b
    start_a = a.get("start", "2000-01-01")[:10]
    end_a = a.get("end", "2000-01-01")[:10]
    start_b = b.get("start", "2000-01-01")[:10]
    end_b = b.get("end", "2000-01-01")[:10]
    from datetime import date as _date
    label_a = f"{_date_label(_date.fromisoformat(start_a))} to {_date_label(_date.fromisoformat(end_a))}"
    label_b = f"{_date_label(_date.fromisoformat(start_b))} to {_date_label(_date.fromisoformat(end_b))}"
    source = f"NEW CASES / FOLLOW-UP CASES | {source_label}"

    def _vert_list(d):
        raw = d.get("vertical") or {}
        return [int(raw.get(v, {}).get("total", 0) or 0) for v in VERTICALS]

    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    cover = _new_slide(prs)
    _text(cover, MARGIN, Inches(2.6), Inches(12.5), Inches(0.6), "Comparative Weekly Wellness Data", 27, color=RED, bold=True, align=PP_ALIGN.CENTER)
    _text(cover, MARGIN, Inches(3.35), Inches(12.5), Inches(0.4), f"{label_a} to {label_b}", 15, color=BLACK, bold=True, align=PP_ALIGN.CENTER)
    _footer(cover, 1, source)
    summary = _new_slide(prs)
    _section(summary, "WEEKLY REPORT SUMMARY", label_a, label_b)
    _table(summary, Inches(0.6), Inches(1.45), Inches(12.1), Inches(2.0), [
        ["Metric", label_a, label_b, "Difference"],
        ["New Cases", a["new"], b["new"], b["new"] - a["new"]],
        ["Follow-up Cases", a["followup"], b["followup"], b["followup"] - a["followup"]],
        ["Grand Total Cases", a["grand"], b["grand"], b["grand"] - a["grand"]],
    ])
    _footer(summary, 2, source)
    details = _new_slide(prs)
    _section(details, "WEEKLY DATA DETAILS", label_a, label_b)
    _footer(details, 3, source)
    from pptx.enum.chart import XL_CHART_TYPE
    va = _vert_list(a)
    vb = _vert_list(b)
    _pair_pies(prs, 4, "VERTICALS — NEW / FOLLOW UP", label_a, label_b, list(VERTICAL_LABELS), va, vb, source)
    _pair_pies(prs, 5, "STAKEHOLDER", label_a, label_b, legacy.STAKEHOLDERS, list(a["stakeholder"].values()), list(b["stakeholder"].values()), source)
    _pair_pies(prs, 6, "RANGE OF CONCERN ADDRESSED", label_a, label_b, legacy.CONCERNS, list(a["concern"].values()), list(b["concern"].values()), source)
    for page, title, labels_list, la, lb in (
        (7, "Comparative data of cases", list(VERTICAL_LABELS), va, vb),
        (8, "Comparative data towards Range of Concern addressed", legacy.CONCERNS, list(a["concern"].values()), list(b["concern"].values())),
        (9, "Comparative data of Stakeholders", legacy.STAKEHOLDERS, list(a["stakeholder"].values()), list(b["stakeholder"].values())),
    ):
        slide = _new_slide(prs)
        _section(slide, title, label_a, label_b)
        _chart(slide, XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.7), CHART_TOP, Inches(11.9), CHART_H,
               labels_list, [(label_a, la), (label_b, lb)])
        _footer(slide, page, source)
    points = _new_slide(prs)
    _text(points, MARGIN, Inches(0.35), Inches(12.5), Inches(0.45), "PROPOSED POINTS FROM WELLNESS CENTRE:", 18, color=RED, bold=True, align=PP_ALIGN.CENTER)
    _footer(points, 10, source)
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


def _build_monthly_fallback(period_a, period_b, insights, source_label):
    """Fallback monthly builder."""
    return _build_weekly_fallback(period_a, period_b, source_label)


def _build_yearly_fallback(periods_a, periods_b, fy1_label, fy2_label):
    """Fallback yearly builder."""
    output = io.BytesIO()
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    _new_slide(prs, "Wellness Data- Report")
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


def _build_normal_week_fallback(period, source_label="uploaded Wellness Excel reports"):
    """Fallback single-week builder."""
    d = _ensure_dict(period) if not isinstance(period, dict) else period
    start = d.get("start", "2000-01-01")[:10]
    end = d.get("end", "2000-01-01")[:10]
    from datetime import date as _date
    label = f"{_date_label(_date.fromisoformat(start))} to {_date_label(_date.fromisoformat(end))}"
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    cover = _new_slide(prs)
    _text(cover, MARGIN, Inches(2.6), Inches(12.5), Inches(0.6),
          "Weekly Wellness Data Report", 27, color=RED, bold=True,
          align=PP_ALIGN.CENTER)
    _text(cover, MARGIN, Inches(3.35), Inches(12.5), Inches(0.4),
          label, 15, color=BLACK, bold=True, align=PP_ALIGN.CENTER)
    _footer(cover, 1, f"Data source: {source_label}")
    summary = _new_slide(prs)
    _section(summary, f"WEEKLY SUMMARY | {label}", "", "")
    _table(summary, Inches(0.6), Inches(1.45), Inches(12.1), Inches(2.0), [
        ["Metric", "Count", "% of Total"],
        ["New Cases", d["new"], f"{d['new']/d['grand']*100:.1f}%" if d['grand'] else "0%"],
        ["Follow-up Cases", d["followup"], f"{d['followup']/d['grand']*100:.1f}%" if d['grand'] else "0%"],
        ["Grand Total", d["grand"], "100%"],
    ])
    _footer(summary, 2, f"Data source: {source_label}")
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


def _build_normal_monthly_fallback(period, source_label="uploaded Wellness Excel reports"):
    """Fallback single-month builder."""
    return _build_normal_week_fallback(period, source_label)


def _build_normal_yearly_fallback(periods, source_label="uploaded Wellness Excel reports"):
    """Fallback single-year builder using first period dict for a basic report."""
    if isinstance(periods, (list, tuple)) and periods:
        period = periods[0]
    else:
        period = periods
    return _build_normal_week_fallback(period, source_label)
