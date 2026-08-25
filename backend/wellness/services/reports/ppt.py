"""Wellness Centre PPT Generation Engine — SAME DESIGN, BETTER ALIGNMENT.

Preserves the EXISTING visual identity from the reference PPTs:
  • Same donut/circle charts where originals used pies
  • Same bar charts where originals used bars
  • Same color palette: #6366f1 #f59e0b #10b981 #ef4444 #8b5cf6 #0ea5e9 #f43f5e #14b8a6
  • Same modules, same structure, same terminology

ADDED: Mandatory NEW CASES / FOLLOW-UP CASES footer on every data slide.

Monthly: 13 slides   Weekly: 10 slides
All 4 verticals: WLN Ctr (WC) · Team A (TA) · Your Dost (YD) · Myndwell (MW)
All 5 data groups: Gender · Mode · Referral · Concern · Stakeholder
"""

from __future__ import annotations

import io
import textwrap
import warnings
from datetime import date
from typing import Any

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from wellness.models import CaseRow, Period

# ═══════════════════════════════ CONSTANTS ═══════════════════════════════════

VERTICALS       = ["WC", "TA", "YD", "MW"]
VERTICAL_LABELS = ["WLN Ctr (WC)", "Team A", "Your Dost (YD)", "Myndwell (MW)"]

GENDERS       = ["Male", "Female", "Others / Not to Say"]
GENDER_FIELDS = ["gender_male", "gender_female", "gender_other"]

MODES         = ["Online", "In-Person", "Phone"]
MODE_FIELDS   = ["mode_online", "mode_in_person", "mode_phone"]

REFERRALS     = ["Self", "Director / Kushal Calls",
                 "Dean / HoD / Faculty / Insti Hosp",
                 "Friend / Family", "Mitr / Saathi"]
REFERRAL_LABELS_SHORT = ["Self", "Director", "Dean/HoD", "Friend/Family", "Mitr/Saathi"]
REFERRAL_FIELDS = ["referral_self", "referral_director", "referral_dean",
                   "referral_friend", "referral_mitr"]

CONCERNS = [
    "Anxiety/Depresn/Panic/OCD", "Acute Stress/Trauma", "Career/Acad",
    "Inter-personal", "Self-Devlp", "Clinical", "Addiction",
    "Medical/Health Issues", "Suicidal Ideation/Self-harm",
]
CONCERN_LABELS_SHORT = CONCERNS  # already short
CONCERN_FIELDS = [
    "concern_anxiety", "concern_stress", "concern_career", "concern_interpersonal",
    "concern_self_dev", "concern_clinical", "concern_addiction", "concern_medical",
    "concern_suicidal",
]

STAKEHOLDERS = ["UG", "PG", "Ph.D.", "Dual Degree", "IIT Faculty/Staff",
                "Employee Family", "Post Doc/Proj Asso", "Not Able to Identify"]
STAKE_LABELS_SHORT = STAKEHOLDERS
STAKE_FIELDS = [
    "stake_ug", "stake_pg", "stake_phd", "stake_dual", "stake_faculty",
    "stake_employee_family", "stake_postdoc", "stake_unidentified",
]

# ═══════════════════════ EXISTING PPT COLOR PALETTE ═════════════════════════
# Extracted from the original reference PPT — DO NOT CHANGE c new = #6366f1, fu = #f59e0b, grand = #10b981 

_PALETTE = ["#4472C4", "#ED7D31", "#A5A5A5", "#FFC000",
            "#5B9BD5", "#70AD47", "#264478", "#9B59B6", "#7F7F7F"]

C_NEW       = "#4472C4"
C_FU        = "#ED7D31"
C_GRAND     = "#70AD47"

# Header/footer use the same palette family
C_HEADER_BG  = "#4472C4"
C_FOOTER_BG  = "#FFFFFF"
C_FOOTER_BDR = "#D9D9D9"
C_TITLE_TXT  = "#FF0000"
C_BODY_TXT   = "#000000"

# ═══════════════════════════ SLIDE DIMENSIONS ═══════════════════════════════

SLIDE_W  = Inches(13.33)
SLIDE_H  = Inches(7.5)
HEADER_H = Inches(0.95)
FOOTER_H = Inches(0.68)
MARGIN   = Inches(0.35)
CONTENT_TOP = Inches(1.05)
CONTENT_W   = SLIDE_W - 2 * MARGIN

DPI = 180

# ═══════════════════════════ DATA HELPERS ═══════════════════════════════════


class _ZeroRow:
    """Stand-in for a missing CaseRow — all numeric fields return 0."""
    def __getattr__(self, item):
        if item.startswith("_"):
            raise AttributeError(item)
        return 0


def index_rows(period: Period) -> dict:
    d = {(r.case_type, r.vertical): r for r in period.case_rows.all()}
    for ct in ("new", "followup"):
        for v in VERTICALS:
            if (ct, v) not in d:
                d[(ct, v)] = _ZeroRow()
    return d


def _case_sum(rows, field, case_type):
    return sum(getattr(rows[(case_type, v)], field) for v in VERTICALS)


def _total_sum(rows, field):
    return _case_sum(rows, field, "new") + _case_sum(rows, field, "followup")


def _new_total(rows):
    return sum(getattr(rows[("new", v)], "total_cases") for v in VERTICALS)


def _fu_total(rows):
    return sum(getattr(rows[("followup", v)], "total_cases") for v in VERTICALS)


def _grand_total(rows):
    return _new_total(rows) + _fu_total(rows)


# ═══════════════════════════ DATE HELPERS ═══════════════════════════════════

def _ordinal(day):
    if 10 <= day % 100 <= 20:
        s = "th"
    else:
        s = {1: "st", 2: "nd", 3: "rd"}.get(day % 10, "th")
    return f"{day}{s}"


def ordinal(day):
    return _ordinal(day)


def fmt_short_range(a: date, b: date) -> str:
    if a.month == b.month and a.year == b.year:
        return f"{_ordinal(a.day)}–{_ordinal(b.day)} {b.strftime('%B %Y')}"
    return f"{_ordinal(a.day)} {a.strftime('%B')} – {_ordinal(b.day)} {b.strftime('%B %Y')}"


def fmt_full_range(a: date, b: date) -> str:
    if a.month == b.month and a.year == b.year:
        return f"{_ordinal(a.day)} to {_ordinal(b.day)} {b.strftime('%B %Y')}"
    return f"{_ordinal(a.day)} {a.strftime('%B %Y')} to {_ordinal(b.day)} {b.strftime('%B %Y')}"


# ═══════════════════════ PPTX LAYOUT HELPERS ════════════════════════════════

def _hex(c):
    c = c.lstrip("#")
    return RGBColor(int(c[:2], 16), int(c[2:4], 16), int(c[4:6], 16))


def _prs_init():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, left, top, w, h, fill, line=None):
    sh = slide.shapes.add_shape(1, left, top, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = _hex(fill)
    if line:
        sh.line.color.rgb = _hex(line)
        sh.line.width = Pt(0.5)
    else:
        sh.line.fill.background()
    return sh


def _txt(slide, text, left, top, w, h, sz=Pt(12), bold=False, color="#000",
         align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = sz
    r.font.bold = bold
    r.font.color.rgb = _hex(color)
    r.font.name = font
    return tb


def _img(slide, img_bytes, left, top, w, h):
    return slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, w, h)


# ─── Header: thin accent bar + title ───────────────────────────────────────

def _add_header(slide, title, period_str, subtitle=""):
    # Thin accent bar at top
    _rect(slide, 0, 0, SLIDE_W, Inches(0.06), C_HEADER_BG)
    # Title
    _txt(slide, title, MARGIN, Inches(0.15), Inches(9), Inches(0.5),
         sz=Pt(22), bold=True, color=C_TITLE_TXT)
    # Period (right side)
    _txt(slide, period_str, Inches(9.5), Inches(0.18), Inches(3.5), Inches(0.45),
         sz=Pt(12), bold=False, color=C_BODY_TXT, align=PP_ALIGN.RIGHT)
    # Subtitle
    if subtitle:
        _txt(slide, subtitle, MARGIN, Inches(0.6), Inches(9), Inches(0.3),
             sz=Pt(10), color="#64748b")
    # Thin line under header
    _rect(slide, MARGIN, HEADER_H - Inches(0.02), CONTENT_W, Inches(0.02), "#e2e8f0")


# ─── Footer: NEW CASES | FOLLOW-UP CASES ──────────────────────────────────

def _add_footer(slide, new_n, fu_n, grand=None,
                wk_a_lbl=None, wk_b_lbl=None,
                wk_a_new=None, wk_a_fu=None, wk_b_new=None, wk_b_fu=None):
    ft = SLIDE_H - FOOTER_H
    _rect(slide, 0, ft, SLIDE_W, FOOTER_H, C_FOOTER_BG, line=C_FOOTER_BDR)
    # Top border accent
    _rect(slide, 0, ft, SLIDE_W, Inches(0.03), C_HEADER_BG)

    if wk_a_lbl and wk_b_lbl:
        # Comparative weekly footer — two halves
        half = SLIDE_W // 2
        _rect(slide, half - Pt(0.5), ft, Pt(1), FOOTER_H, C_FOOTER_BDR)
        for i, (lbl, n, f) in enumerate([
            (wk_a_lbl, wk_a_new, wk_a_fu),
            (wk_b_lbl, wk_b_new, wk_b_fu),
        ]):
            x = Inches(0.4) + half * i
            _txt(slide, lbl, x, ft + Inches(0.06), Inches(5.8), Inches(0.18),
                 sz=Pt(7.5), bold=True, color="#374151")
            _txt(slide, f"NEW CASES: {n}   |   FOLLOW-UP: {f}",
                 x, ft + Inches(0.28), Inches(5.8), Inches(0.3),
                 sz=Pt(12), bold=True, color=C_TITLE_TXT)
        return

    # Standard footer
    sections = 3 if grand is not None else 2
    col_w = SLIDE_W / sections
    items = [("NEW CASES", str(new_n), C_NEW),
             ("FOLLOW-UP CASES", str(fu_n), C_FU)]
    if grand is not None:
        items.append(("GRAND TOTAL", str(grand), C_GRAND))

    for i, (lbl, val, clr) in enumerate(items):
        cx = col_w * i
        if i > 0:
            _rect(slide, cx, ft, Pt(1), FOOTER_H, C_FOOTER_BDR)
        # Color accent dot
        _rect(slide, cx + Inches(0.3), ft + Inches(0.15),
              Inches(0.08), Inches(0.08), clr)
        _txt(slide, lbl, cx + Inches(0.45), ft + Inches(0.08),
             col_w - Inches(0.5), Inches(0.2),
             sz=Pt(8), bold=True, color="#64748b")
        _txt(slide, val, cx + Inches(0.45), ft + Inches(0.28),
             col_w - Inches(0.5), Inches(0.3),
             sz=Pt(18), bold=True, color=C_TITLE_TXT)


# ═══════════════════════ CHART RENDERERS ════════════════════════════════════

def _fig_bytes(fig):
    buf = io.BytesIO()
    with warnings.catch_warnings():
        warnings.simplefilter("ignore")
        fig.savefig(buf, format="png", facecolor="white", dpi=DPI, bbox_inches="tight")
    plt.close(fig)
    return buf.getvalue()


def _style_ax(ax):
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#e2e8f0")
    ax.spines["bottom"].set_color("#e2e8f0")
    ax.tick_params(colors="#475569", labelsize=8.5)
    ax.yaxis.grid(True, color="#f1f5f9", linewidth=0.8, zorder=0)
    ax.set_axisbelow(True)


# ─── DONUT CHART (preserves original circle/pie style) ─────────────────────

def _donut(values, labels, title, center_text="", colors=None,
           figsize=(5, 4.5), show_legend=True):
    """Single donut chart — same style as existing PPT circular charts."""
    fig, ax = plt.subplots(figsize=figsize, facecolor="white")
    ax.set_aspect("equal")

    # Filter zeros
    nz = [(v, l, i) for i, (v, l) in enumerate(zip(values, labels)) if v > 0]
    if not nz:
        ax.text(0.5, 0.5, "No Data", ha="center", va="center",
                fontsize=14, transform=ax.transAxes, color="#94a3b8")
        ax.axis("off")
        return _fig_bytes(fig)

    vals, lbls, idxs = zip(*nz)
    clrs = colors or _PALETTE
    wedge_colors = [clrs[i % len(clrs)] for i in idxs]

    wedges, texts, autotexts = ax.pie(
        vals, labels=None,
        autopct=lambda p: f"{p:.0f}%" if p >= 4 else "",
        colors=wedge_colors, startangle=90,
        pctdistance=0.78,
        wedgeprops={"width": 0.38, "edgecolor": "white", "linewidth": 2},
        textprops={"fontsize": 8, "color": "white", "fontweight": "bold"},
    )
    for at in autotexts:
        at.set_fontsize(8)
        at.set_fontweight("bold")
        at.set_color("white")

    # Center text (total)
    if center_text:
        ax.text(0, 0.06, str(center_text), ha="center", va="center",
                fontsize=22, fontweight="bold", color=C_TITLE_TXT)
        ax.text(0, -0.14, "Total", ha="center", va="center",
                fontsize=8, color="#94a3b8")

    ax.set_title(title, fontsize=11, fontweight="bold", color=C_TITLE_TXT,
                 pad=10, loc="center")
    if show_legend:
        leg_labels = [f"{l}  ({v})" for l, v in zip(lbls, vals)]
        ax.legend(wedges, leg_labels, loc="center left", bbox_to_anchor=(1.02, 0.5),
                  fontsize=8.5, frameon=False, labelspacing=0.8)
    fig.tight_layout(pad=0.5)
    return _fig_bytes(fig)


def _paired_donuts(vals_a, vals_b, labels, title_a, title_b,
                   suptitle="", colors=None, figsize=(12.5, 4.8)):
    """Two donut charts side-by-side — used for Period A vs B or New vs Follow-up."""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=figsize, facecolor="white")
    if suptitle:
        fig.suptitle(suptitle, fontsize=12, fontweight="bold", color=C_TITLE_TXT, y=0.98)

    for ax, vals, title in [(ax1, vals_a, title_a), (ax2, vals_b, title_b)]:
        ax.set_aspect("equal")
        nz = [(v, l, i) for i, (v, l) in enumerate(zip(vals, labels)) if v > 0]
        if not nz:
            ax.text(0.5, 0.5, "No Data", ha="center", va="center",
                    fontsize=12, color="#94a3b8", transform=ax.transAxes)
            ax.axis("off")
            ax.set_title(title, fontsize=10, fontweight="bold", color=C_TITLE_TXT)
            continue

        vs, ls, idxs = zip(*nz)
        clrs = colors or _PALETTE
        wc = [clrs[i % len(clrs)] for i in idxs]
        total = sum(vs)

        wedges, _, autotexts = ax.pie(
            vs, labels=None,
            autopct=lambda p: f"{p:.0f}%" if p >= 4 else "",
            colors=wc, startangle=90, pctdistance=0.78,
            wedgeprops={"width": 0.38, "edgecolor": "white", "linewidth": 2},
            textprops={"fontsize": 8, "color": "white", "fontweight": "bold"},
        )
        for at in autotexts:
            at.set_fontsize(7.5)
            at.set_fontweight("bold")
            at.set_color("white")

        # Center total
        ax.text(0, 0.06, str(total), ha="center", va="center",
                fontsize=18, fontweight="bold", color=C_TITLE_TXT)
        ax.text(0, -0.14, "Total", ha="center", va="center",
                fontsize=7.5, color="#94a3b8")
        ax.set_title(title, fontsize=10, fontweight="bold", color=C_TITLE_TXT, pad=8)

    # Shared legend
    all_labels = [f"{l}" for l in labels]
    clrs = colors or _PALETTE
    handles = [mpatches.Patch(facecolor=clrs[i % len(clrs)], edgecolor="white",
                              label=all_labels[i]) for i in range(len(labels))]
    fig.legend(handles=handles, loc="lower center", ncol=min(len(labels), 5),
               fontsize=8.5, frameon=False, bbox_to_anchor=(0.5, -0.02))
    fig.tight_layout(pad=0.8, rect=[0, 0.06, 1, 0.94])
    return _fig_bytes(fig)


# ─── BAR CHART (for Referral, Concern, Stakeholder monthly) ───────────────

def _hbar(labels, series, title, figsize=(12.5, 4.5)):
    """Horizontal bar chart — same style as existing PPT bar charts."""
    n = len(labels)
    ns = len(series)
    fig, ax = plt.subplots(figsize=figsize, facecolor="white")
    w = 0.7 / ns
    y = np.arange(n)

    for i, (name, vals, clr) in enumerate(series):
        offset = (i - (ns - 1) / 2) * w
        bars = ax.barh(y + offset, vals, height=w * 0.9, label=name,
                       color=clr, zorder=3, edgecolor="white", linewidth=0.5)
        for bar, v in zip(bars, vals):
            if v > 0:
                ax.text(bar.get_width() + 0.15, bar.get_y() + bar.get_height() / 2,
                        str(v), ha="left", va="center", fontsize=7.5,
                        color="#374151", fontweight="bold")

    ax.set_yticks(y)
    ax.set_yticklabels([textwrap.fill(str(l), 20) for l in labels], fontsize=8.5)
    ax.invert_yaxis()
    ax.set_title(title, fontsize=11, fontweight="bold", color=C_TITLE_TXT,
                 loc="left", pad=8)
    if ns > 1:
        ax.legend(frameon=False, fontsize=9, loc="lower right")
    _style_ax(ax)
    ax.spines["left"].set_visible(False)
    ax.tick_params(left=False)
    fig.tight_layout(pad=0.6)
    return _fig_bytes(fig)


def _clustered_bar(labels, series, title, figsize=(12.5, 4.2)):
    """Vertical clustered bar chart — same style as existing PPT bar charts."""
    n = len(labels)
    ns = len(series)
    fig, ax = plt.subplots(figsize=figsize, facecolor="white")
    w = 0.7 / ns
    x = np.arange(n)

    for i, (name, vals, clr) in enumerate(series):
        offset = (i - (ns - 1) / 2) * w
        bars = ax.bar(x + offset, vals, width=w * 0.9, label=name,
                      color=clr, zorder=3, edgecolor="white", linewidth=0.5)
        for bar, v in zip(bars, vals):
            if v > 0:
                ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.2,
                        str(v), ha="center", va="bottom", fontsize=7.5,
                        color="#374151", fontweight="bold")

    ax.set_xticks(x)
    ax.set_xticklabels([textwrap.fill(str(l), 12) for l in labels], fontsize=8.5)
    ax.set_title(title, fontsize=11, fontweight="bold", color=C_TITLE_TXT,
                 loc="left", pad=8)
    ax.legend(frameon=False, fontsize=9, loc="upper right")
    _style_ax(ax)
    fig.tight_layout(pad=0.6)
    return _fig_bytes(fig)


# ─── SUMMARY CARDS ─────────────────────────────────────────────────────────

def _summary_cards(new_n, fu_n, grand, period, figsize=(12.5, 2.8)):
    """Three summary cards matching existing PPT visual style."""
    fig, axes = plt.subplots(1, 3, figsize=figsize, facecolor="white")
    cards = [
        ("NEW CASES", new_n, "#eef2ff", C_NEW),
        ("FOLLOW-UP CASES", fu_n, "#fffbeb", C_FU),
        ("GRAND TOTAL", grand, "#ecfdf5", C_GRAND),
    ]
    for ax, (label, val, bg, accent) in zip(axes, cards):
        ax.set_xlim(0, 1)
        ax.set_ylim(0, 1)
        ax.set_facecolor(bg)
        ax.axis("off")
        # Top accent bar
        ax.add_patch(mpatches.Rectangle((0, 0.88), 1, 0.12, color=accent,
                     transform=ax.transAxes, clip_on=False))
        ax.text(0.5, 0.94, label, ha="center", va="center",
                fontsize=10, fontweight="bold", color="white",
                transform=ax.transAxes)
        ax.text(0.5, 0.48, str(val), ha="center", va="center",
                fontsize=42, fontweight="bold", color=accent,
                transform=ax.transAxes)
        ax.text(0.5, 0.12, period, ha="center", va="center",
                fontsize=8.5, color="#94a3b8", transform=ax.transAxes)
        for spine in ax.spines.values():
            spine.set_edgecolor("#e2e8f0")
            spine.set_linewidth(1)

    fig.tight_layout(pad=0.6)
    return _fig_bytes(fig)


# ─── DATA TABLE ────────────────────────────────────────────────────────────

def _table_img(headers, rows, title="", figsize=(12.5, 3.0)):
    fig, ax = plt.subplots(figsize=figsize, facecolor="white")
    ax.axis("off")
    if title:
        ax.set_title(title, fontsize=10, fontweight="bold", color=C_TITLE_TXT,
                     loc="left", pad=6)

    nw = len(headers)
    col_widths = [1.0 / nw] * nw
    cell_text = [[str(c) for c in row] for row in rows]
    table = ax.table(cellText=cell_text, colLabels=headers, loc="center",
                     cellLoc="center", colWidths=col_widths)
    table.auto_set_font_size(False)
    table.set_fontsize(8.5)
    table.scale(1, 1.5)
    for (r, c), cell in table.get_celld().items():
        cell.set_edgecolor("#e2e8f0")
        if r == 0:
            cell.set_facecolor(C_HEADER_BG)
            cell.set_text_props(color="white", fontweight="bold")
        elif r % 2 == 0:
            cell.set_facecolor("#f8fafc")
    try:
        fig.tight_layout(pad=0.3)
    except Exception:
        pass
    return _fig_bytes(fig)


# ═══════════════════════ MONTHLY BUILDER (12 slides) ════════════════════════

def _resolve_bullets(period, previous=None, insights=None):
    if insights and insights.get("insights"):
        return [b["text"] for b in insights["insights"][:6]]
    from wellness.services.insights import analyze_period
    return [b["text"] for b in analyze_period(period, previous)["insights"][:6]]


def _ai_slide(prs, period_str, bullets, bar_labels, bar_series, pie_values, pie_labels):
    """AI DATA INSIGHTS slide — one pie chart + one bar chart + insight bullets."""
    s = _blank(prs)
    _add_header(s, "AI DATA INSIGHTS", period_str)
    _img(s, _donut(pie_values, pie_labels, "Concern Mix", figsize=(5, 4.5)),
         Inches(0.3), Inches(1.35), Inches(5.4), Inches(3.7))
    _img(s, _clustered_bar(bar_labels, bar_series, "Cases by Vertical", figsize=(7.6, 4.0)),
         Inches(5.9), Inches(1.35), Inches(7.1), Inches(3.7))
    bullets_text = "\n".join(f"•  {b}" for b in bullets)
    _txt(s, bullets_text, MARGIN, Inches(5.15), CONTENT_W, Inches(1.6),
         sz=Pt(11), color=C_BODY_TXT)
    return s


def build_monthly(period: Period, insights=None) -> bytes:
    rows = index_rows(period)
    prs = _prs_init()

    new_n = _new_total(rows)
    fu_n  = _fu_total(rows)
    grand = _grand_total(rows)
    mo = period.period_end.strftime("%B %Y").upper()

    def _slide(title, sub=""):
        s = _blank(prs)
        _add_header(s, title, mo, sub)
        return s

    def _ft(s, show_grand=False):
        _add_footer(s, new_n, fu_n, grand if show_grand else None)

    # ── S1: Title + Summary Cards ──────────────────────────────────────────
    s1 = _slide("MONTHLY WELLNESS DATA", mo)
    _img(s1, _summary_cards(new_n, fu_n, grand, mo),
         MARGIN, CONTENT_TOP, CONTENT_W, Inches(2.6))
    # Small New vs FU clustered bar below cards
    _img(s1, _clustered_bar(
        VERTICAL_LABELS,
        [("New Cases", [getattr(rows[("new", v)], "total_cases") for v in VERTICALS], C_NEW),
         ("Follow-up", [getattr(rows[("followup", v)], "total_cases") for v in VERTICALS], C_FU)],
        f"Cases by Vertical — {mo}", figsize=(12.5, 2.3)),
         MARGIN, CONTENT_TOP + Inches(2.7), CONTENT_W, Inches(2.3))
    _ft(s1, show_grand=True)

    # ── S2: Vertical Distribution (DONUT — total by vertical) ──────────────
    s2 = _slide("VERTICAL DISTRIBUTION")
    vert_totals = [getattr(rows[("new", v)], "total_cases") +
                   getattr(rows[("followup", v)], "total_cases") for v in VERTICALS]
    _img(s2, _donut(vert_totals, VERTICAL_LABELS,
                    f"Total Cases by Vertical — {mo}", center_text=str(grand),
                    figsize=(6, 5)),
         Inches(0.3), CONTENT_TOP, Inches(6.5), Inches(4.9))
    # Table on the right
    _img(s2, _table_img(
        ["Vertical", "New", "Follow-up", "Total"],
        [[VERTICAL_LABELS[i],
          getattr(rows[("new", v)], "total_cases"),
          getattr(rows[("followup", v)], "total_cases"),
          vert_totals[i]] for i, v in enumerate(VERTICALS)]
        + [["ALL", new_n, fu_n, grand]],
        figsize=(5.8, 2.5)),
         Inches(6.9), CONTENT_TOP + Inches(0.5), Inches(6.1), Inches(2.4))
    _ft(s2)

    # ── S3: New vs Follow-up by Vertical (PAIRED DONUTS) ───────────────────
    s3 = _slide("NEW VS FOLLOW-UP BY VERTICAL")
    new_by_v = [getattr(rows[("new", v)], "total_cases") for v in VERTICALS]
    fu_by_v  = [getattr(rows[("followup", v)], "total_cases") for v in VERTICALS]
    _img(s3, _paired_donuts(new_by_v, fu_by_v, VERTICAL_LABELS,
                            f"NEW CASES ({new_n})", f"FOLLOW-UP CASES ({fu_n})",
                            suptitle=f"New vs Follow-up by Vertical — {mo}",
                            figsize=(12.5, 4.5)),
         MARGIN, CONTENT_TOP, CONTENT_W, Inches(4.8))
    _ft(s3)

    # ── S4: Gender Distribution (PAIRED DONUTS — New / Follow-up) ──────────
    s4 = _slide("GENDER DISTRIBUTION")
    g_new = [_case_sum(rows, f, "new") for f in GENDER_FIELDS]
    g_fu  = [_case_sum(rows, f, "followup") for f in GENDER_FIELDS]
    _img(s4, _paired_donuts(g_new, g_fu, GENDERS,
                            f"NEW CASES ({sum(g_new)})", f"FOLLOW-UP ({sum(g_fu)})",
                            suptitle=f"Gender Distribution — {mo}",
                            figsize=(12.5, 4.5)),
         MARGIN, CONTENT_TOP, CONTENT_W, Inches(4.8))
    _ft(s4)

    # ── S5: Mode of Session (PAIRED DONUTS) ────────────────────────────────
    s5 = _slide("MODE OF SESSION")
    m_new = [_case_sum(rows, f, "new") for f in MODE_FIELDS]
    m_fu  = [_case_sum(rows, f, "followup") for f in MODE_FIELDS]
    _img(s5, _paired_donuts(m_new, m_fu, MODES,
                            f"NEW CASES ({sum(m_new)})", f"FOLLOW-UP ({sum(m_fu)})",
                            suptitle=f"Mode of Session — {mo}",
                            figsize=(12.5, 4.5)),
         MARGIN, CONTENT_TOP, CONTENT_W, Inches(4.8))
    _ft(s5)

    # ── S6: Referral Type (HORIZONTAL BAR — same as original) ──────────────
    s6 = _slide("REFERRAL TYPE")
    r_new = [_case_sum(rows, f, "new") for f in REFERRAL_FIELDS]
    r_fu  = [_case_sum(rows, f, "followup") for f in REFERRAL_FIELDS]
    _img(s6, _hbar(REFERRALS,
                   [("New Cases", r_new, C_NEW), ("Follow-up", r_fu, C_FU)],
                   f"Referral Sources — {mo}", figsize=(12.5, 4.0)),
         MARGIN, CONTENT_TOP, CONTENT_W, Inches(4.2))
    _img(s6, _table_img(
        ["Referral", "New", "FU", "Total"],
        [[REFERRAL_LABELS_SHORT[i], r_new[i], r_fu[i], r_new[i]+r_fu[i]]
         for i in range(5)], figsize=(12.5, 1.1)),
         MARGIN, CONTENT_TOP + Inches(4.3), CONTENT_W, Inches(1.0))
    _ft(s6)

    # ── S7: Stakeholder (HORIZONTAL BAR — same as original) ────────────────
    s7 = _slide("STAKEHOLDER")
    st_new = [_case_sum(rows, f, "new") for f in STAKE_FIELDS]
    st_fu  = [_case_sum(rows, f, "followup") for f in STAKE_FIELDS]
    _img(s7, _hbar(STAKEHOLDERS,
                   [("New Cases", st_new, C_NEW), ("Follow-up", st_fu, C_FU)],
                   f"Stakeholder — {mo}", figsize=(12.5, 4.3)),
         MARGIN, CONTENT_TOP, CONTENT_W, Inches(4.5))
    _img(s7, _table_img(
        ["Stakeholder", "New", "FU", "Total"],
        [[STAKE_LABELS_SHORT[i], st_new[i], st_fu[i], st_new[i]+st_fu[i]]
         for i in range(8)], figsize=(12.5, 1.6)),
         MARGIN, CONTENT_TOP + Inches(4.55), CONTENT_W, Inches(1.1))
    _ft(s7)

    # ── S8: Range of Concern (HORIZONTAL BAR — same as original) ───────────
    s8 = _slide("RANGE OF CONCERN ADDRESSED")
    c_new = [_case_sum(rows, f, "new") for f in CONCERN_FIELDS]
    c_fu  = [_case_sum(rows, f, "followup") for f in CONCERN_FIELDS]
    _img(s8, _hbar(CONCERNS,
                   [("New Cases", c_new, C_NEW), ("Follow-up", c_fu, C_FU)],
                   f"Range of Concern — {mo}", figsize=(12.5, 4.5)),
         MARGIN, CONTENT_TOP, CONTENT_W, Inches(4.6))
    _ft(s8)

    # ── S9: Summary Table ──────────────────────────────────────────────────
    s9 = _slide("WELLNESS DATA SUMMARY")
    def _vt(field, v):
        return sum(getattr(rows[(ct, v)], field) for ct in ("new", "followup"))
    def _s(field):
        return sum(_vt(field, v) for v in VERTICALS)
    def _gstr(v):
        return f"{_vt('gender_male',v)}/{_vt('gender_female',v)}/{_vt('gender_other',v)}"
    def _mstr(v):
        return f"{_vt('mode_online',v)}/{_vt('mode_in_person',v)}/{_vt('mode_phone',v)}"

    tbl_rows = [
        ["Grand Total"] + [vert_totals[i] for i in range(4)] + [grand],
        ["New / Follow-up"] + [f"{getattr(rows[('new',v)],'total_cases')}/{getattr(rows[('followup',v)],'total_cases')}"
                               for v in VERTICALS] + [f"{new_n}/{fu_n}"],
        ["Gender (M/F/O)"] + [_gstr(v) for v in VERTICALS] + [
            f"{_s('gender_male')}/{_s('gender_female')}/{_s('gender_other')}"],
        ["Mode (On/IP/Ph)"] + [_mstr(v) for v in VERTICALS] + [
            f"{_s('mode_online')}/{_s('mode_in_person')}/{_s('mode_phone')}"],
    ]
    _img(s9, _table_img(
        ["Metric"] + VERTICAL_LABELS + ["Total"], tbl_rows,
        f"WELLNESS DATA SUMMARY — {mo}", figsize=(12.5, 4.0)),
         MARGIN, CONTENT_TOP, CONTENT_W, Inches(4.0))
    _ft(s9)

    # ── S10: Concerns New vs Follow-up (BAR — same as original) ────────────
    s10 = _slide("RANGE OF CONCERN — NEW VS FOLLOW-UP")
    _img(s10, _hbar(CONCERNS,
                    [("New", c_new, C_NEW), ("Follow-up", c_fu, C_FU)],
                    f"Range of Concern (New vs Follow-up) — {mo}",
                    figsize=(12.5, 5.0)),
         MARGIN, CONTENT_TOP, CONTENT_W, Inches(5.1))
    _ft(s10)

    # ── S11: Grand Total ───────────────────────────────────────────────────
    s11 = _slide("GRAND TOTAL")
    _img(s11, _clustered_bar(
        VERTICAL_LABELS,
        [("New", [getattr(rows[("new", v)], "total_cases") for v in VERTICALS], C_NEW),
         ("Follow-up", [getattr(rows[("followup", v)], "total_cases") for v in VERTICALS], C_FU)],
        f"Grand Total by Vertical — {mo}", figsize=(12.5, 3.0)),
         MARGIN, CONTENT_TOP, CONTENT_W, Inches(3.2))
    _img(s11, _table_img(
        ["Vertical", "New", "Follow-up", "Grand Total"],
        [[VERTICAL_LABELS[i], getattr(rows[("new",v)],"total_cases"),
          getattr(rows[("followup",v)],"total_cases"), vert_totals[i]]
         for i, v in enumerate(VERTICALS)]
        + [["ALL VERTICALS", new_n, fu_n, grand]],
        figsize=(12.5, 2.2)),
         MARGIN, CONTENT_TOP + Inches(3.3), CONTENT_W, Inches(2.1))
    _ft(s11, show_grand=True)

    # ── S12: AI Data Insights ──────────────────────────────────────────────
    s12 = _ai_slide(
        prs, mo, _resolve_bullets(period, insights=insights),
        VERTICAL_LABELS,
        [("New Cases", [getattr(rows[("new", v)], "total_cases") for v in VERTICALS], C_NEW),
         ("Follow-up", [getattr(rows[("followup", v)], "total_cases") for v in VERTICALS], C_FU)],
        [_case_sum(rows, f, "new") + _case_sum(rows, f, "followup") for f in CONCERN_FIELDS],
        CONCERN_LABELS_SHORT,
    )
    _ft(s12, show_grand=True)

    # ── S13: Activities (no footer) ────────────────────────────────────────
    s13 = _blank(prs)
    _add_header(s13, "ACTIVITIES & PROGRAMMES", mo)
    _txt(s13, "Please add programme / activity / proposed activities content here.",
         MARGIN, CONTENT_TOP, CONTENT_W, Inches(5),
         sz=Pt(12), color=C_BODY_TXT)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ═══════════════════════ AI COMPARISON BUILDER (5 slides) ════════════════════
# Used by the "Data Analysis AI" tab for week-over-week, month-over-month and
# year-over-year comparisons. Takes an insights dict from
# wellness.services.insights.compare_periods.

def _aggregate_year_periods(periods):
    """Aggregate monthly CaseRows into the yearly report's three verticals."""
    fields = [
        "total_cases", *GENDER_FIELDS, *MODE_FIELDS, *REFERRAL_FIELDS,
        *CONCERN_FIELDS, *STAKE_FIELDS,
    ]
    result = {
        (case_type, vertical): {field: 0 for field in fields}
        for case_type in ("new", "followup")
        for vertical in ("WC", "TA", "YD", "MW")
    }
    for period in periods:
        rows = index_rows(period)
        for case_type in ("new", "followup"):
            for vertical in ("WC", "TA", "YD", "MW"):
                source = rows[(case_type, vertical)]
                target = result[(case_type, vertical)]
                for field in fields:
                    target[field] += int(getattr(source, field, 0) or 0)
    return result


def build_yearly(periods_a, periods_b,
                 fy1_label="FY 2024-25", fy2_label="FY 2025-26") -> bytes:
    """Build the reference 12-slide annual comparison deck from monthly periods."""
    from wellness.services.reports.reference_ppt import build_yearly as build_reference_yearly
    return build_reference_yearly(periods_a, periods_b, fy1_label, fy2_label)


def build_annual(periods, insights=None) -> bytes:
    """Merged Jan-Dec annual analysis deck (single year, no comparison)."""
    from wellness.services.reports.reference_ppt import build_normal_yearly as build_reference_annual
    return build_reference_annual(periods, insights=insights)


def build_weekly_comparison(period_a, period_b, source_label="uploaded Wellness Excel reports") -> bytes:
    from wellness.services.reports.reference_ppt import build_weekly_comparison as build_reference_weekly
    return build_reference_weekly(period_a, period_b, source_label)


def build_monthly_comparison(period_a, period_b, insights=None, source_label="uploaded Wellness Excel reports") -> bytes:
    from wellness.services.reports.reference_ppt import build_monthly_comparison as build_reference_monthly
    return build_reference_monthly(period_a, period_b, insights, source_label)


def build_ai_comparison(period_a, period_b, insights=None, insert_into_ppt=True, comparison_type="week") -> bytes:
    from wellness.services.insights import compare_periods

    if insights is None:
        insights = compare_periods(period_a, period_b, comparison_type)

    a = index_rows(period_a)
    b = index_rows(period_b)
    prs = _prs_init()

    a_new, a_fu = _new_total(a), _fu_total(a)
    b_new, b_fu = _new_total(b), _fu_total(b)
    a_grand, b_grand = a_new + a_fu, b_new + b_fu

    sa = fmt_short_range(period_a.period_start, period_a.period_end)
    sb = fmt_short_range(period_b.period_start, period_b.period_end)
    type_label = insights.get("comparison_label") or "Period-over-Period"
    period_str = f"{type_label}:  {sa}  vs  {sb}"

    def _slide(title, sub=""):
        s = _blank(prs)
        _add_header(s, title, period_str, sub)
        return s

    def _ft(s):
        _add_footer(s, a_new, a_fu,
                    wk_a_lbl=sa, wk_b_lbl=sb,
                    wk_a_new=a_new, wk_a_fu=a_fu,
                    wk_b_new=b_new, wk_b_fu=b_fu)

    # ── S1: Title + two summary card rows ─────────────────────────────────
    s1 = _slide("AI DATA ANALYSIS")
    _img(s1, _summary_cards(a_new, a_fu, a_grand, sa, figsize=(12.5, 2.0)),
         MARGIN, CONTENT_TOP, CONTENT_W, Inches(2.0))
    _txt(s1, f"{sa.upper()}  (baseline)", MARGIN, CONTENT_TOP - Inches(0.16), CONTENT_W, Inches(0.16),
         sz=Pt(8), bold=True, color=C_NEW)
    _img(s1, _summary_cards(b_new, b_fu, b_grand, sb, figsize=(12.5, 2.0)),
         MARGIN, CONTENT_TOP + Inches(2.5), CONTENT_W, Inches(2.0))
    _txt(s1, f"{sb.upper()}  (current)", MARGIN, CONTENT_TOP + Inches(2.34), CONTENT_W, Inches(0.16),
         sz=Pt(8), bold=True, color=C_FU)
    _ft(s1)

    # ── S2: Key changes (clustered bar + change table) ─────────────────────
    s2 = _slide("KEY CHANGES")
    _img(s2, _clustered_bar(
        VERTICAL_LABELS,
        [(f"{sa} New", [getattr(a[('new', v)], 'total_cases') for v in VERTICALS], C_NEW),
         (f"{sa} FU",  [getattr(a[('followup', v)], 'total_cases') for v in VERTICALS], "#a5b4fc"),
         (f"{sb} New", [getattr(b[('new', v)], 'total_cases') for v in VERTICALS], C_FU),
         (f"{sb} FU",  [getattr(b[('followup', v)], 'total_cases') for v in VERTICALS], "#fcd34d")],
        "New & Follow-up by Vertical — Both Periods", figsize=(12.5, 3.6)),
         MARGIN, CONTENT_TOP, CONTENT_W, Inches(3.7))
    t = insights["totals"]
    _img(s2, _table_img(
        ["Metric", sa, sb, "Change", "% Change"],
        [["Total Cases", a_grand, b_grand, b_grand - a_grand,
          f"{t['pct_total']:+.1f}%" if t["pct_total"] is not None else "—"],
         ["New", a_new, b_new, b_new - a_new,
          f"{t['pct_new']:+.1f}%" if t["pct_new"] is not None else "—"],
         ["Follow-up", a_fu, b_fu, b_fu - a_fu,
          f"{t['pct_followup']:+.1f}%" if t["pct_followup"] is not None else "—"]],
        figsize=(12.5, 1.3)),
         MARGIN, CONTENT_TOP + Inches(3.8), CONTENT_W, Inches(1.2))
    _ft(s2)

    # ── S3: Vertical distribution (paired donuts) ──────────────────────────
    s3 = _slide("VERTICAL DISTRIBUTION")
    va = [getattr(a[("new", v)], "total_cases") + getattr(a[("followup", v)], "total_cases")
          for v in VERTICALS]
    vb = [getattr(b[("new", v)], "total_cases") + getattr(b[("followup", v)], "total_cases")
          for v in VERTICALS]
    _img(s3, _paired_donuts(va, vb, VERTICAL_LABELS, sa, sb,
                            suptitle="Total Cases by Vertical", figsize=(12.5, 4.8)),
         MARGIN, CONTENT_TOP, CONTENT_W, Inches(5.0))
    _ft(s3)

    # ── S4: Top category changes (table) ───────────────────────────────────
    s4 = _slide("TOP CATEGORY CHANGES")
    movers = insights.get("movers") or []
    m_rows = [[m["category"].capitalize(), m["label"], f"{m['a']} → {m['b']}",
               f"{'▲' if m['delta'] > 0 else '▼'} {'+' if m['delta'] >= 0 else ''}{m['delta']}"]
              for m in movers[:14]]
    if not m_rows:
        m_rows = [["—", "No category-level shifts detected", "—", "—"]]
    _img(s4, _table_img(
        ["Category", "Dimension", f"{sa} → {sb}", "Change"],
        m_rows, "Largest Category Shifts", figsize=(12.5, 5.2)),
         MARGIN, CONTENT_TOP, CONTENT_W, Inches(5.2))
    _ft(s4)

    # ── S5: AI insights (bullets + concern donut for current period) ───────
    s5 = _slide("AI INSIGHTS")
    if insert_into_ppt:
        bullets = [b["text"] for b in (insights.get("insights") or [])][:8]
        bullets_text = "\n".join(f"•  {t}" for t in bullets) or "No insights available."
    else:
        bullets_text = "AI analysis was not inserted into this PPT (insert_into_ppt=False)."
    _img(s5, _donut([getattr(b[("new", v)], "total_cases") + getattr(b[("followup", v)], "total_cases")
                     for v in VERTICALS], VERTICAL_LABELS,
                    "Cases by Vertical — Current Period", center_text=str(b_grand),
                    figsize=(5.4, 4.6)),
         Inches(0.3), CONTENT_TOP, Inches(5.6), Inches(4.7))
    _txt(s5, bullets_text, Inches(6.1), CONTENT_TOP, Inches(6.9), Inches(4.7),
         sz=Pt(11), color=C_BODY_TXT)
    _ft(s5)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ═══════════════════════ WEEKLY BUILDER (9 slides) ══════════════════════════

def build_weekly(period_a: Period, period_b: Period, insights=None) -> bytes:
    from wellness.services.reports.reference_ppt import build_weekly_comparison as build_reference_weekly
    return build_reference_weekly(period_a, period_b)
