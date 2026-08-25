"""Tests for the AI period-over-period comparison (week / month / year)."""

import io

import pytest
from pptx import Presentation
from rest_framework.test import APIClient

from tests.builders import fields
from tests.test_reports import _save
from wellness.models import User
from wellness.services.insights import compare_periods
from wellness.services.reports import exports, ppt


@pytest.fixture
def admin_user(db):
    return User.objects.create_user(username="cmp1", password="x", role=User.Role.ADMIN)


def _row(total=0, m=0, f=0, **kw):
    d = fields()
    d.update({"total_cases": total, "gender_male": m, "gender_female": f, "gender_other": 0})
    d.update(kw)
    return d


@pytest.fixture
def weekly_pair(db, admin_user):
    a = _save(
        admin_user, "weekly", "15th July", "21st July 2026",
        {"WLN Ctr": _row(10, 6, 4, concern_anxiety=6, stake_ug=10, referral_self=10, mode_online=10),
         "Team A": _row(), "Your Dost": _row(), "Myndwell": _row()},
        {"WLN Ctr": _row(20, 10, 10), "Team A": _row(), "Your Dost": _row(), "Myndwell": _row()},
        secondary={"total_sessions": {"Total": 90}, "no_show_turn_up": {"Total": 10}},
    )
    b = _save(
        admin_user, "weekly", "22nd July", "28th July 2026",
        {"WLN Ctr": _row(20, 12, 8, concern_stress=14, concern_anxiety=6, stake_pg=20, referral_mitr=20),
         "Team A": _row(), "Your Dost": _row(), "Myndwell": _row()},
        {"WLN Ctr": _row(25, 15, 10), "Team A": _row(), "Your Dost": _row(), "Myndwell": _row()},
        secondary={"total_sessions": {"Total": 120}, "no_show_turn_up": {"Total": 40}},
    )
    return a, b


@pytest.fixture
def monthly_pair(db, admin_user):
    a = _save(
        admin_user, "monthly", "01st July", "31st July 2026",
        {"WLN Ctr": _row(30, 18, 12, concern_anxiety=20), "Team A": _row(),
         "Your Dost": _row(20, 10, 10), "Myndwell": _row()},
        {"WLN Ctr": _row(60, 35, 25), "Team A": _row(), "Your Dost": _row(), "Myndwell": _row()},
    )
    b = _save(
        admin_user, "monthly", "01st August", "31st August 2026",
        {"WLN Ctr": _row(45, 25, 20, concern_anxiety=40), "Team A": _row(),
         "Your Dost": _row(25, 15, 10), "Myndwell": _row()},
        {"WLN Ctr": _row(70, 40, 30), "Team A": _row(), "Your Dost": _row(), "Myndwell": _row()},
    )
    return a, b


@pytest.fixture
def yearly_pair(db, admin_user):
    a = _save(
        admin_user, "monthly", "01st July", "31st July 2025",
        {"WLN Ctr": _row(25, 15, 10), "Team A": _row(), "Your Dost": _row(), "Myndwell": _row()},
        {"WLN Ctr": _row(40, 22, 18), "Team A": _row(), "Your Dost": _row(), "Myndwell": _row()},
    )
    b = _save(
        admin_user, "monthly", "01st July", "31st July 2026",
        {"WLN Ctr": _row(45, 25, 20), "Team A": _row(), "Your Dost": _row(), "Myndwell": _row()},
        {"WLN Ctr": _row(70, 40, 30), "Team A": _row(), "Your Dost": _row(), "Myndwell": _row()},
    )
    return a, b


class TestCompareService:
    def test_week_over_week(self, weekly_pair):
        a, b = weekly_pair
        r = compare_periods(a, b, "week")
        assert r["comparison_label"] == "Week-over-Week"
        assert r["totals"]["delta_total"] == 15
        assert r["totals"]["delta_sessions"] == 30
        assert r["movers"][0]["delta"] == 20  # stake_pg 0 -> 20 largest
        assert any("total cases rose by 15" in i["text"] for i in r["insights"])

    def test_month_over_month(self, monthly_pair):
        a, b = monthly_pair
        r = compare_periods(a, b, "month")
        assert r["comparison_label"] == "Month-over-Month"
        assert r["totals"]["delta_total"] == 30
        assert r["category_deltas"]["concern"]
        assert any("total cases rose by 30" in i["text"] for i in r["insights"])

    def test_year_over_year(self, yearly_pair):
        a, b = yearly_pair
        r = compare_periods(a, b, "year")
        assert r["comparison_label"] == "Year-over-Year"
        assert r["totals"]["delta_total"] == 50
        assert any("total cases rose by 50" in i["text"] for i in r["insights"])

    def test_baseline_ordering(self, weekly_pair):
        # Passing the later period first must still compare earlier -> later.
        a, b = weekly_pair
        r = compare_periods(b, a, "week")
        assert r["period_a"]["id"] == a.id
        assert r["period_b"]["id"] == b.id

    def test_no_data(self, db, admin_user):
        from tests.test_reports import _row as rw
        a = _save(admin_user, "weekly", "15th July", "21st July 2026",
                  {t: fields() for t in ("WLN Ctr", "Team A", "Your Dost", "Myndwell")},
                  {t: fields() for t in ("WLN Ctr", "Team A", "Your Dost", "Myndwell")})
        b = _save(admin_user, "weekly", "22nd July", "28th July 2026",
                  {t: fields() for t in ("WLN Ctr", "Team A", "Your Dost", "Myndwell")},
                  {t: fields() for t in ("WLN Ctr", "Team A", "Your Dost", "Myndwell")})
        r = compare_periods(a, b, "week")
        assert r["insights"][0]["text"] == "No cases were recorded in either period."


class TestComparisonAPI:
    def _client(self, admin_user):
        c = APIClient()
        c.force_authenticate(user=admin_user)
        return c

    def test_weekly_endpoint(self, admin_user, weekly_pair):
        a, b = weekly_pair
        r = self._client(admin_user).get(f"/api/insights/compare?type=week&from_id={a.id}&to_id={b.id}")
        assert r.status_code == 200
        assert r.data["comparison_label"] == "Week-over-Week"
        assert r.data["totals"]["delta_total"] == 15

    def test_monthly_endpoint(self, admin_user, monthly_pair):
        a, b = monthly_pair
        r = self._client(admin_user).get(f"/api/insights/compare?type=month&from_id={a.id}&to_id={b.id}")
        assert r.status_code == 200
        assert r.data["comparison_label"] == "Month-over-Month"

    def test_yearly_endpoint(self, admin_user, yearly_pair):
        a, b = yearly_pair
        r = self._client(admin_user).get(f"/api/insights/compare?type=year&from_id={a.id}&to_id={b.id}")
        assert r.status_code == 200
        assert r.data["comparison_label"] == "Year-over-Year"

    def test_type_mismatch_rejected(self, admin_user, weekly_pair, monthly_pair):
        w, _ = weekly_pair
        m, _ = monthly_pair
        c = self._client(admin_user)
        assert c.get(f"/api/insights/compare?type=week&from_id={w.id}&to_id={m.id}").status_code == 400

    def test_same_year_rejected(self, admin_user, monthly_pair):
        a, b = monthly_pair
        r = self._client(admin_user).get(f"/api/insights/compare?type=year&from_id={a.id}&to_id={b.id}")
        assert r.status_code == 400

    def test_missing_period(self, admin_user, weekly_pair):
        a, _ = weekly_pair
        r = self._client(admin_user).get(f"/api/insights/compare?type=week&from_id={a.id}&to_id=99999")
        assert r.status_code == 400

    def test_requires_auth(self, weekly_pair):
        a, b = weekly_pair
        r = APIClient().get(f"/api/insights/compare?type=week&from_id={a.id}&to_id={b.id}")
        assert r.status_code == 401


class TestComparisonExports:
    def _client(self, admin_user):
        c = APIClient()
        c.force_authenticate(user=admin_user)
        return c

    def test_ppt_builder(self, weekly_pair):
        a, b = weekly_pair
        r = compare_periods(a, b, "week")
        buf = ppt.build_ai_comparison(a, b, r)
        assert buf[:2] == b"PK"
        prs = Presentation(io.BytesIO(buf))
        assert len(list(prs.slides)) == 5
        all_text = " ".join(
            sh.text_frame.text for slide in prs.slides for sh in slide.shapes if sh.has_text_frame
        )
        assert "WEEK-OVER-WEEK" in all_text or "Week-over-Week" in all_text
        assert "AI DATA ANALYSIS" in all_text

    def test_excel_builder(self, weekly_pair):
        a, b = weekly_pair
        r = compare_periods(a, b, "week")
        data = exports.build_comparison_excel(a, b, r)
        assert data[:2] == b"PK"
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(data))
        assert "AI Comparison" in wb.sheetnames
        ws = wb["AI Comparison"]
        assert ws.cell(row=1, column=1).value == "AI DATA ANALYSIS"
        texts = {ws.cell(row=row, column=1).value for row in range(1, 30)}
        assert "Total cases" in texts
        assert "Category breakdown — period-over-period" in texts

    def test_ppt_via_api(self, admin_user, weekly_pair):
        a, b = weekly_pair
        r = self._client(admin_user).post(
            "/api/reports/generate",
            {"format": "comparison_ppt", "compare_type": "week", "from_id": a.id, "to_id": b.id},
            format="json",
        )
        assert r.status_code == 200
        assert b"".join(r.streaming_content)[:2] == b"PK"

    def test_monthly_reference_ppt_via_api(self, admin_user, monthly_pair):
        a, b = monthly_pair
        r = self._client(admin_user).post(
            "/api/reports/generate",
            {"format": "comparison_ppt", "compare_type": "month", "from_id": a.id, "to_id": b.id},
            format="json",
        )
        assert r.status_code == 200
        deck = b"".join(r.streaming_content)
        prs = Presentation(io.BytesIO(deck))
        assert len(prs.slides) == 13
        text = " ".join(
            shape.text_frame.text
            for slide in prs.slides
            for shape in slide.shapes
            if shape.has_text_frame
        )
        assert "VERTICALS TOTAL" in text
        assert "Team A" in text

    def test_excel_via_api(self, admin_user, monthly_pair):
        a, b = monthly_pair
        r = self._client(admin_user).post(
            "/api/reports/generate",
            {"format": "comparison_xlsx", "compare_type": "month", "from_id": a.id, "to_id": b.id},
            format="json",
        )
        assert r.status_code == 200
        assert b"".join(r.streaming_content)[:2] == b"PK"

    def test_bad_compare_type_via_api(self, admin_user, weekly_pair):
        a, b = weekly_pair
        r = self._client(admin_user).post(
            "/api/reports/generate",
            {"format": "comparison_ppt", "compare_type": "daily", "from_id": a.id, "to_id": b.id},
            format="json",
        )
        assert r.status_code == 400
