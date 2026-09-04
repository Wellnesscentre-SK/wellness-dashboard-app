"""
Reusable slide components for Wellness Centre PPT reports.
Matches Excel Report sheet design: blue headers, orange/green/gray KPI cards,
pie charts with percentage labels, bar charts for New vs Follow-up.
"""

import os
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu
try:
    from ppt_generator import config as C
except ImportError:
    import config as C


def create_presentation():
    prs = Presentation()
    prs.slide_width = C.SLIDE_W
    prs.slide_height = C.SLIDE_H
    return prs


def add_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = C.WHITE
    return slide


def add_logo(slide, left=None, top=None, width=None, height=None):
    if not os.path.exists(C.LOGO_PATH):
        return None
    return slide.shapes.add_picture(
        C.LOGO_PATH,
        left or C.LOGO_LEFT,
        top or C.LOGO_TOP,
        width or C.LOGO_WIDTH,
        height or C.LOGO_HEIGHT,
    )


def add_text(slide, text, left, top, width, height,
             font_size=C.BODY_FONT, bold=False, color=C.BLACK,
             alignment=PP_ALIGN.LEFT, font_name=C.FONT_NAME, word_wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = str(text)
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return box


def add_multiline(slide, lines, left, top, width, height,
                  font_size=C.SECTION_FONT, bold=True, color=C.PRIMARY,
                  alignment=PP_ALIGN.CENTER, line_spacing=1.2):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        p.space_before = Pt(0)
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = str(line)
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = C.FONT_NAME
    return box


def add_section_header(slide, text, left, top, width, height,
                       bg_color=C.PRIMARY, font_color=C.WHITE,
                       font_size=C.SECTION_FONT):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(text)
    run.font.size = font_size
    run.font.bold = True
    run.font.color.rgb = font_color
    run.font.name = C.FONT_NAME
    return shape


def add_section_label(slide, left, top, width, height, lines):
    main_text = lines[0] if lines else ""
    subtitle = " | ".join(lines[1:]) if len(lines) > 1 else ""
    add_section_header(slide, main_text, left, top, width, Inches(0.45), bg_color=C.PRIMARY)
    if subtitle:
        add_text(slide, subtitle, left, top + Inches(0.48), width, Inches(0.25),
                 font_size=C.SMALL_FONT, bold=False, color=C.GRAY, alignment=PP_ALIGN.CENTER)


def add_table(slide, left, top, width, height, data, col_widths=None, header_color=None):
    n_rows = len(data)
    n_cols = len(data[0]) if data else 0
    if n_rows == 0 or n_cols == 0:
        return None
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    hdr_color = header_color or C.TABLE_HEADER_BG
    for r, row_data in enumerate(data):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = C.TABLE_FONT
                paragraph.font.name = C.FONT_NAME
                paragraph.alignment = PP_ALIGN.CENTER
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = C.WHITE
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = hdr_color
                else:
                    paragraph.font.color.rgb = C.BLACK
                    if r % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = C.TABLE_ALT_ROW_BG
                    else:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = C.WHITE
    return table_shape


def add_kpi_cards(slide, cards, top=Inches(0.80)):
    """Add KPI summary cards matching Excel design (colored left bar, value, label).
    cards: list of (label, value, color) tuples.
    """
    n = len(cards)
    if n == 0:
        return
    total_w = 12.6
    gap = 0.25
    card_w = (total_w - gap * (n - 1)) / n
    start_x = (13.33 - total_w) / 2
    for i, (label, value, color) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        # Card background
        shape = slide.shapes.add_shape(1, Inches(x), top, Inches(card_w), Inches(0.85))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        # Value text
        add_text(slide, str(value), Inches(x), top + Inches(0.08),
                 Inches(card_w), Inches(0.45),
                 font_size=Pt(24), bold=True, color=C.WHITE, alignment=PP_ALIGN.CENTER)
        # Label text
        add_text(slide, label, Inches(x), top + Inches(0.50),
                 Inches(card_w), Inches(0.30),
                 font_size=C.SMALL_FONT, bold=False, color=C.WHITE, alignment=PP_ALIGN.CENTER)


def add_kpi_cards_comparison(slide, label_a, val_a, label_b, val_b, top=Inches(0.80)):
    """Add two-period comparison KPI cards."""
    cards = [
        (f"Total ({label_a})", val_a, C.TOTAL_CLR),
        (f"Total ({label_b})", val_b, C.NEW_CASE_CLR),
        ("Difference", val_b - val_a, C.FOLLOWUP_CLR),
    ]
    add_kpi_cards(slide, cards, top)


def add_pie_chart(slide, left, top, width, height,
                  labels, values, colors=None, show_legend=False):
    if not labels or not values or all(v == 0 for v in values):
        add_text(slide, "No data available", left, top, width, height,
                 font_size=C.HEADING_FONT, bold=True, color=C.GRAY,
                 alignment=PP_ALIGN.CENTER)
        return None
    colors = colors or C.PIE_COLORS
    data = CategoryChartData()
    data.categories = [str(l) for l in labels]
    data.add_series("Cases", [int(v or 0) for v in values])
    frame = slide.shapes.add_chart(XL_CHART_TYPE.PIE, left, top, width, height, data)
    chart = frame.chart
    chart.has_legend = show_legend
    if show_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = C.LEGEND_FONT
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_category_name = False
    dl.show_value = False
    dl.show_percentage = True
    dl.number_format = "0%"
    dl.number_format_is_linked = False
    dl.font.size = C.LABEL_FONT
    dl.font.bold = True
    dl.font.color.rgb = C.BLACK
    series = chart.series[0]
    for i, point in enumerate(series.points):
        fill = point.format.fill
        fill.solid()
        fill.fore_color.rgb = colors[i % len(colors)]
    return frame


def add_pie_with_table(slide, left_table, left_chart, top, table_w, chart_w, height,
                       table_data, chart_labels, chart_values, chart_colors=None,
                       title_text=None):
    """Add a table on the left and pie chart on the right, matching Excel layout."""
    if title_text:
        add_text(slide, title_text, left_table, top - Inches(0.35), table_w + chart_w + Inches(0.3),
                 Inches(0.30), font_size=C.HEADING_FONT, bold=True, color=C.PRIMARY,
                 alignment=PP_ALIGN.LEFT)
    add_table(slide, left_table, top, table_w, height, table_data)
    add_pie_chart(slide, left_chart, top, chart_w, height,
                  chart_labels, chart_values, chart_colors)


def add_pie_pair(slide, left_labels, right_labels,
                 left_values, right_values,
                 left_title_lines, right_title_lines, colors=None):
    add_section_label(slide, C.LABEL_LEFT_LEFT, C.LABEL_LEFT_TOP,
                      C.LABEL_LEFT_WIDTH, C.LABEL_LEFT_HEIGHT, left_title_lines)
    add_section_label(slide, C.LABEL_RIGHT_LEFT, C.LABEL_RIGHT_TOP,
                      C.LABEL_RIGHT_WIDTH, C.LABEL_RIGHT_HEIGHT, right_title_lines)
    add_pie_chart(slide, C.PIE_HALF_LEFT, C.PIE_HALF_TOP,
                  C.PIE_HALF_WIDTH, C.PIE_HALF_HEIGHT,
                  left_labels, left_values, colors)
    add_pie_chart(slide, C.PIE_HALF_RIGHT_LEFT, C.PIE_HALF_RIGHT_TOP,
                  C.PIE_HALF_RIGHT_WIDTH, C.PIE_HALF_RIGHT_HEIGHT,
                  right_labels, right_values, colors)


def add_column_chart(slide, left, top, width, height,
                     labels, series_data, legend_position="bottom",
                     title_text=None, series_colors=None, show_percentage=True):
    series_colors = series_colors or C.SERIES_COLORS
    data = CategoryChartData()
    data.categories = [str(l) for l in labels]
    for name, values in series_data:
        data.add_series(name, [int(v or 0) for v in values])
    frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, data)
    chart = frame.chart
    if title_text:
        chart.has_title = True
        chart.chart_title.text_frame.text = title_text
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
    else:
        chart.has_title = False
    chart.has_legend = True
    if legend_position == "right":
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
    else:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = C.LEGEND_FONT
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_value = True
    dl.show_category_name = False
    dl.show_percentage = show_percentage
    dl.number_format = "0" if not show_percentage else '0" ("0%")'
    dl.font.size = C.LABEL_FONT
    dl.font.bold = True
    dl.font.color.rgb = C.BLACK
    for i, series_obj in enumerate(chart.series):
        series_obj.format.fill.solid()
        series_obj.format.fill.fore_color.rgb = series_colors[i % len(series_colors)]
    return frame


def add_bar_chart(slide, left, top, width, height,
                  labels, series_data, legend_position="bottom", series_colors=None):
    series_colors = series_colors or C.SERIES_COLORS
    data = CategoryChartData()
    data.categories = [str(l) for l in labels]
    for name, values in series_data:
        data.add_series(name, [int(v or 0) for v in values])
    frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, data)
    chart = frame.chart
    chart.has_title = False
    chart.has_legend = True
    if legend_position == "right":
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
    else:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = C.LEGEND_FONT
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_value = True
    dl.show_percentage = True
    dl.number_format = '0" ("0%")'
    dl.font.size = C.LABEL_FONT
    dl.font.bold = True
    for i, s in enumerate(chart.series):
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = series_colors[i % len(series_colors)]
    return frame


def add_cover(prs, title, subtitle=None):
    slide = add_slide(prs)
    add_logo(slide)
    add_text(slide, title, Inches(1.5), Inches(2.5), Inches(10.33), Inches(1.0),
             font_size=C.COVER_FONT, bold=True, color=C.PRIMARY, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, subtitle, Inches(1.5), Inches(3.5), Inches(10.33), Inches(0.6),
                 font_size=Pt(18), bold=False, color=C.DARK_GRAY, alignment=PP_ALIGN.CENTER)
    shape = slide.shapes.add_shape(1, Inches(0), Inches(7.0), Inches(13.33), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C.PRIMARY
    shape.line.fill.background()
    add_text(slide, "IIT Madras Wellness Centre", Inches(0.5), Inches(7.05),
             Inches(12.33), Inches(0.4), font_size=C.SMALL_FONT, bold=False,
             color=C.WHITE, alignment=PP_ALIGN.CENTER)
    return slide


def add_proposed_points(prs, points=None):
    slide = add_slide(prs)
    add_section_header(slide, "PROPOSED POINTS FROM WELLNESS CENTRE",
                       Inches(0.15), Inches(0.15), Inches(13.0), Inches(0.55),
                       bg_color=C.PRIMARY)
    if points:
        content = "\n\n".join(f"{i+1}. {p}" for i, p in enumerate(points))
    else:
        content = "[Enter proposed points here]"
    add_text(slide, content, Inches(0.5), Inches(0.90), Inches(12.33), Inches(6.20),
             font_size=C.BODY_FONT, bold=False, color=C.BLACK, alignment=PP_ALIGN.LEFT)
    return slide


def add_wellness_title(slide):
    add_text(slide, "Wellness Data- Report", Inches(4.6), Inches(0.03),
             Inches(4.0), Inches(0.5), font_size=C.COVER_FONT, bold=True,
             color=C.PRIMARY, alignment=PP_ALIGN.CENTER)


def add_full_column(slide, title_text, labels, series_data,
                    legend_position="bottom", series_colors=None,
                    title_font=C.HEADING_FONT, show_percentage=True):
    add_text(slide, title_text, Inches(0.9), Inches(0.51), Inches(11.5), Inches(0.34),
             font_size=title_font, bold=True, color=C.PRIMARY, alignment=PP_ALIGN.CENTER)
    add_column_chart(slide, C.COL_LEFT, C.COL_TOP, C.COL_WIDTH, C.COL_HEIGHT,
                     labels, series_data, legend_position=legend_position,
                     series_colors=series_colors, show_percentage=show_percentage)


def add_data_table_slide(prs, title, headers, rows, page_label=""):
    slide = add_slide(prs)
    add_logo(slide)
    add_section_header(slide, title, Inches(0.15), Inches(0.75),
                       Inches(13.0), Inches(0.45), bg_color=C.PRIMARY)
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_h_val = min(5.5, 0.35 * n_rows)
    table_w_val = min(12.5, 1.5 * n_cols)
    left_val = (13.33 - table_w_val) / 2
    add_table(slide, Inches(left_val), Inches(1.35),
              Inches(table_w_val), Inches(table_h_val), [headers] + rows)
    return slide


def add_insights_slide(prs, title, insights_list, period_label=""):
    slide = add_slide(prs)
    add_logo(slide)
    header_text = f"{title}  |  {period_label}" if period_label else title
    add_section_header(slide, header_text, Inches(0.15), Inches(0.75),
                       Inches(13.0), Inches(0.45), bg_color=C.PRIMARY)
    y = Inches(1.40)
    for i, insight in enumerate(insights_list[:10], 1):
        add_text(slide, f"{i}.", Inches(0.6), y, Inches(0.4), Inches(0.45),
                 font_size=C.BODY_FONT, bold=True, color=C.PRIMARY, alignment=PP_ALIGN.RIGHT)
        add_text(slide, insight, Inches(1.1), y, Inches(11.9), Inches(0.45),
                 font_size=Pt(12), bold=False, color=C.BLACK, alignment=PP_ALIGN.LEFT)
        y += Inches(0.55)
    return slide


def add_metric_card(slide, left, top, width, height, label, value, color=C.PRIMARY):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = C.WHITE
    shape.line.color.rgb = C.LIGHT_GRAY
    shape.line.width = Pt(1)
    bar = slide.shapes.add_shape(1, left, top, width, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text(slide, str(value), left, top + Inches(0.15), width, Inches(0.5),
             font_size=Pt(28), bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text(slide, label, left, top + Inches(0.65), width, Inches(0.3),
             font_size=C.SMALL_FONT, bold=False, color=C.GRAY, alignment=PP_ALIGN.CENTER)


def add_footer(slide, text="IIT Madras Wellness Centre"):
    shape = slide.shapes.add_shape(1, Inches(0), Inches(7.0), Inches(13.33), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C.PRIMARY
    shape.line.fill.background()
    add_text(slide, text, Inches(0.5), Inches(7.05), Inches(12.33), Inches(0.4),
             font_size=C.SMALL_FONT, bold=False, color=C.WHITE, alignment=PP_ALIGN.CENTER)
