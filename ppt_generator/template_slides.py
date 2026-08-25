"""
Template slides — replicates 'Gender distribution_Range of Concern
addressed_Stake holder.pptx'.

Three slides, each with the logo, a rounded-rectangle title bar and a
full-width clustered column chart. Categories repeat once per period
(spacer column between groups), and the series are the verticals
WC / YD / MW so every period shows its vertical-wise split side by side,
exactly like the reference deck.

Usage:
    add_template_slides(prs, period)                      # single period
    add_template_slides(prs, a, b, label_a, label_b)      # two periods
"""

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import config as C
import components as CM

TITLE_BAR_LEFT = Inches(3.37)
TITLE_BAR_TOP = Inches(0.56)
TITLE_BAR_W = Inches(6.60)
TITLE_BAR_H = Inches(0.52)

LABEL_Y = Inches(1.12)
CHART_SINGLE_TOP = Inches(1.30)
CHART_DUAL_TOP = Inches(1.45)
CHART_W = Inches(13.33)
CHART_H = Inches(6.00)

GENDER_CATS = ["Male", "Female", "Others/\nNot to say"]
CONCERN_CATS = [
    "Anxiety/\nDepresn/\nPanic / OCD", "Acute\n Stress / Trauma", "Career /\n Acad,",
    "Inter-\npersonal", "Self-\nDevlp", "Clinical", "Addiction",
    "Medical /\n Health issues", "Suicidal Ideation\n / Self-harm",
]
STAKEHOLDER_CATS = [
    "UG", "PG", "Ph.D", "Dual Degree", "IIT Faculty / Staff",
    "Employee Family", "Post Doc\nProj Asso", "Not Able to Identify",
]

_TITLES = {
    "gender": "GENDER DISTRIBUTION OF CASES",
    "concern": "RANGE OF CONCERN ADDRESSED",
    "stakeholder": "STAKE HOLDER",
}

_DIMS = [
    ("gender", C.GENDER_LABELS, GENDER_CATS),
    ("concern", C.CONCERN_LABELS, CONCERN_CATS),
    ("stakeholder", C.STAKEHOLDER_LABELS, STAKEHOLDER_CATS),
]

# display series name -> source vertical keys merged together (WC absorbs TA)
_VERT_SERIES = [("WC", ("WC", "TA")), ("YD", ("YD",)), ("MW", ("MW",))]
_SERIES_COLORS = [C.WC_COLOR, C.YD_COLOR, C.MW_COLOR]


def _vertical_breakdown(p, dim, key):
    """{series_name: count} for one category of one period."""
    bv = p.get("by_vertical")
    if bv:
        out = {}
        for name, sources in _VERT_SERIES:
            out[name] = sum(int(bv.get(s, {}).get(dim, {}).get(key, 0) or 0)
                            for s in sources)
        return out
    # fallback: spread the category total across verticals by case share
    raw = p.get("vertical", {})

    def _vt(k):
        return int(raw.get(k, {}).get("total", 0) or 0)

    total = int(p.get(dim, {}).get(key, 0) or 0)
    merged = {}
    for name, sources in _VERT_SERIES:
        vals = [_vt(s) for s in sources]
        if not sum(vals):
            vals = [_vt("WLC" if name == "WC" else name)]
        merged[name] = sum(vals)
    grand = sum(merged.values())
    if not grand:
        return {name: 0 for name, _ in _VERT_SERIES}
    return {name: round(total * merged[name] / grand) for name in merged}


def _title_bar(slide, text):
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 TITLE_BAR_LEFT, TITLE_BAR_TOP,
                                 TITLE_BAR_W, TITLE_BAR_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = C.PRIMARY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = C.WHITE
    run.font.name = C.FONT_NAME


def _period_labels(slide, la, lb):
    if la:
        CM.add_text(slide, la, Inches(0.15), LABEL_Y, Inches(6.4), Inches(0.28),
                    font_size=Pt(11), bold=True, color=C.PRIMARY,
                    alignment=PP_ALIGN.CENTER)
    if lb:
        CM.add_text(slide, lb, Inches(6.93), LABEL_Y, Inches(6.4), Inches(0.28),
                    font_size=Pt(11), bold=True, color=C.PRIMARY,
                    alignment=PP_ALIGN.CENTER)


def _chart(slide, cats, series, top):
    data = CategoryChartData()
    data.categories = cats
    for name, values in series:
        data.add_series(name, values)
    frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                   0, top, CHART_W, CHART_H, data)
    chart = frame.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = C.LEGEND_FONT
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.show_value = True
    dl.show_category_name = False
    dl.show_percentage = False
    dl.number_format = "0"
    dl.number_format_is_linked = False
    dl.font.size = Pt(8)
    dl.font.bold = True
    dl.font.color.rgb = C.BLACK
    for i, s in enumerate(chart.series):
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = _SERIES_COLORS[i % len(_SERIES_COLORS)]
    try:
        chart.category_axis.tick_labels.font.size = Pt(8)
        chart.value_axis.tick_labels.font.size = Pt(8)
    except Exception:
        pass
    return frame


def _add_dim_slide(prs, dim, title, keys, cats, pa, pb, la, lb):
    slide = CM.add_slide(prs)
    CM.add_logo(slide)
    _title_bar(slide, title)

    series_acc = {name: [] for name, _ in _VERT_SERIES}
    if pb is None:
        for k in keys:
            bd = _vertical_breakdown(pa, dim, k)
            for name in series_acc:
                series_acc[name].append(bd[name])
        _chart(slide, list(cats),
               [(n, series_acc[n]) for n, _ in _VERT_SERIES], CHART_SINGLE_TOP)
        return slide

    _period_labels(slide, la, lb)
    gap = 0 if dim == "gender" else 1          # template keeps no spacer on gender
    all_cats = list(cats) + [""] * gap + list(cats)
    for idx, p in enumerate((pa, pb)):
        for k in keys:
            bd = _vertical_breakdown(p, dim, k)
            for name in series_acc:
                series_acc[name].append(bd[name])
        if idx == 0 and gap:
            for name in series_acc:
                series_acc[name].append(0)
    _chart(slide, all_cats,
           [(n, series_acc[n]) for n, _ in _VERT_SERIES], CHART_DUAL_TOP)
    return slide


def add_template_slides(prs, period_a, period_b=None,
                        label_a="", label_b=""):
    """Append the three template slides (Gender / Concern / Stake Holder)."""
    for dim, keys, cats in _DIMS:
        _add_dim_slide(prs, dim, _TITLES[dim], keys, cats,
                       period_a, period_b, label_a, label_b)
