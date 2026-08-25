"""
Yearly Wellness Comparison Report — Excel-style layout.
KPI comparison, verticals, concern, stakeholder with tables+pies+bars, insights.
"""

import io
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import config as C
import components as CM
import template_slides as TS


def _aggregate(periods_list: list) -> dict:
    agg = {
        "new": 0, "followup": 0, "grand": 0,
        "gender": {}, "mode": {}, "referral": {},
        "concern": {}, "stakeholder": {},
        "vertical": {v: {"new": 0, "followup": 0, "total": 0} for v in C.VERT_KEYS},
        "by_vertical": {},
    }
    for p in periods_list:
        agg["new"] += p.get("new", 0)
        agg["followup"] += p.get("followup", 0)
        agg["grand"] += p.get("grand", 0)
        for dim in ("gender", "mode", "referral", "concern", "stakeholder"):
            for k, v in p.get(dim, {}).items():
                agg[dim][k] = agg[dim].get(k, 0) + v
        raw_vert = p.get("vertical", {})
        combined = C.combine_verticals(raw_vert)
        for vk in C.VERT_KEYS:
            pv = combined.get(vk, {})
            for sub in ("new", "followup", "total"):
                agg["vertical"][vk][sub] += pv.get(sub, 0)
        for src, dims in (p.get("by_vertical") or {}).items():
            tgt = agg["by_vertical"].setdefault(src, {})
            for dim, cats in dims.items():
                t = tgt.setdefault(dim, {})
                for k, v in cats.items():
                    t[k] = t.get(k, 0) + v
    return agg


def build(periods_a: list, periods_b: list,
          fy1_label: str = "FY 2024-25", fy2_label: str = "FY 2025-26",
          fy1_period_label: str = None, fy2_period_label: str = None,
          proposed_points: list = None) -> bytes:
    prs = CM.create_presentation()
    a = _aggregate(periods_a)
    b = _aggregate(periods_b)
    v_labels = [C.VERTICALS[v] for v in C.VERT_KEYS]
    vk = C.VERT_KEYS
    st_labels = C.STAKEHOLDER_LABELS
    c_labels = C.CONCERN_LABELS

    # ── SLIDE 1: COVER ────────────────────────────────────────────────────────
    CM.add_cover(prs, "Yearly Wellness Data",
                 f"Comparative Analysis: {fy1_label} vs. {fy2_label}")

    # ── SLIDE 2: KPI COMPARISON + SUMMARY TABLES ──────────────────────────────
    s2 = CM.add_slide(prs)
    CM.add_logo(s2)
    CM.add_kpi_cards(s2, [
        (f"Total ({fy1_label})", a["grand"], C.TOTAL_CLR),
        (f"Total ({fy2_label})", b["grand"], C.NEW_CASE_CLR),
        ("Difference", b["grand"] - a["grand"], C.FOLLOWUP_CLR),
        ("New Cases", f"{a['new']} → {b['new']}", C.TOP_VERT_CLR),
    ])
    CM.add_table(s2, Inches(0.45), Inches(1.85), Inches(12.4), Inches(2.0), [
        ["Metric", fy1_label, fy2_label, "Difference", "% Change"],
        ["Total Cases", a["grand"], b["grand"], b["grand"] - a["grand"],
         _pct(a["grand"], b["grand"])],
        ["New Cases", a["new"], b["new"], b["new"] - a["new"],
         _pct(a["new"], b["new"])],
        ["Follow-up Cases", a["followup"], b["followup"],
         b["followup"] - a["followup"], _pct(a["followup"], b["followup"])],
    ])
    CM.add_table(s2, Inches(0.45), Inches(4.00), Inches(12.4), Inches(2.5), [
        ["Vertical", fy1_label, fy2_label, "Difference"],
        *[[v_labels[i], a["vertical"][v]["total"], b["vertical"][v]["total"],
           b["vertical"][v]["total"] - a["vertical"][v]["total"]]
          for i, v in enumerate(vk)],
    ])

    # ── SLIDE 3: VERTICALS PIE + TABLE (FY1) ──────────────────────────────────
    s3 = CM.add_slide(prs)
    CM.add_logo(s3)
    a_vals = [a["vertical"][v]["total"] for v in vk]
    CM.add_pie_chart(s3, Inches(0.15), Inches(1.35), Inches(7.0), Inches(5.5),
                     v_labels, a_vals)
    CM.add_table(s3, Inches(7.5), Inches(1.60), Inches(5.5), Inches(3.0), [
        ["Vertical", "New", "Follow-up", "Total"],
        *[[v_labels[i], a["vertical"][v]["new"],
           a["vertical"][v]["followup"], a["vertical"][v]["total"]]
          for i, v in enumerate(vk)],
    ])
    CM.add_text(s3, f"VERTICALS  |  {fy1_label}", Inches(0.15), Inches(0.95),
                Inches(13.0), Inches(0.30), font_size=C.HEADING_FONT, bold=True,
                color=C.PRIMARY, alignment=PP_ALIGN.CENTER)

    # ── SLIDE 4: VERTICALS PIE + TABLE (FY2) ──────────────────────────────────
    s4 = CM.add_slide(prs)
    CM.add_logo(s4)
    b_vals = [b["vertical"][v]["total"] for v in vk]
    CM.add_pie_chart(s4, Inches(0.15), Inches(1.35), Inches(7.0), Inches(5.5),
                     v_labels, b_vals)
    CM.add_table(s4, Inches(7.5), Inches(1.60), Inches(5.5), Inches(3.0), [
        ["Vertical", "New", "Follow-up", "Total"],
        *[[v_labels[i], b["vertical"][v]["new"],
           b["vertical"][v]["followup"], b["vertical"][v]["total"]]
          for i, v in enumerate(vk)],
    ])
    CM.add_text(s4, f"VERTICALS  |  {fy2_label}", Inches(0.15), Inches(0.95),
                Inches(13.0), Inches(0.30), font_size=C.HEADING_FONT, bold=True,
                color=C.PRIMARY, alignment=PP_ALIGN.CENTER)

    # ── SLIDE 5: STAKEHOLDER PIE (FY1) ────────────────────────────────────────
    s5 = CM.add_slide(prs)
    CM.add_logo(s5)
    CM.add_section_header(s5, f"STAKEHOLDER  |  {fy1_label}", Inches(0.15), Inches(0.75),
                          Inches(13.0), Inches(0.45), bg_color=C.PRIMARY)
    CM.add_pie_chart(s5, Inches(0.3), Inches(1.35), Inches(12.7), Inches(5.5),
                     st_labels, [a["stakeholder"].get(s, 0) for s in st_labels])

    # ── SLIDE 6: STAKEHOLDER PIE (FY2) ────────────────────────────────────────
    s6 = CM.add_slide(prs)
    CM.add_logo(s6)
    CM.add_section_header(s6, f"STAKEHOLDER  |  {fy2_label}", Inches(0.15), Inches(0.75),
                          Inches(13.0), Inches(0.45), bg_color=C.PRIMARY)
    CM.add_pie_chart(s6, Inches(0.3), Inches(1.35), Inches(12.7), Inches(5.5),
                     st_labels, [b["stakeholder"].get(s, 0) for s in st_labels])

    # ── SLIDE 7: CONCERN PIE (FY1) ────────────────────────────────────────────
    s7 = CM.add_slide(prs)
    CM.add_logo(s7)
    CM.add_section_header(s7, f"RANGE OF CONCERN  |  {fy1_label}", Inches(0.15), Inches(0.75),
                          Inches(13.0), Inches(0.45), bg_color=C.PRIMARY)
    CM.add_pie_chart(s7, Inches(0.3), Inches(1.35), Inches(12.7), Inches(5.5),
                     c_labels, [a["concern"].get(c, 0) for c in c_labels])

    # ── SLIDE 8: CONCERN PIE (FY2) ────────────────────────────────────────────
    s8 = CM.add_slide(prs)
    CM.add_logo(s8)
    CM.add_section_header(s8, f"RANGE OF CONCERN  |  {fy2_label}", Inches(0.15), Inches(0.75),
                          Inches(13.0), Inches(0.45), bg_color=C.PRIMARY)
    CM.add_pie_chart(s8, Inches(0.3), Inches(1.35), Inches(12.7), Inches(5.5),
                     c_labels, [b["concern"].get(c, 0) for c in c_labels])

    # ── SLIDE 9-11: COMPARISON BAR CHARTS ─────────────────────────────────────
    for title, labels, vals_a, vals_b in [
        ("VERTICALS COMPARISON", v_labels,
         [a["vertical"][v]["total"] for v in vk], [b["vertical"][v]["total"] for v in vk]),
        ("STAKEHOLDER COMPARISON", st_labels,
         [a["stakeholder"].get(s, 0) for s in st_labels], [b["stakeholder"].get(s, 0) for s in st_labels]),
        ("CONCERN COMPARISON", c_labels,
         [a["concern"].get(c, 0) for c in c_labels], [b["concern"].get(c, 0) for c in c_labels]),
    ]:
        s = CM.add_slide(prs)
        CM.add_logo(s)
        CM.add_full_column(s, f"{title}  |  {fy1_label} vs {fy2_label}",
                           labels, [(fy1_label, vals_a), (fy2_label, vals_b)],
                           legend_position="bottom")

    # ── TEMPLATE SLIDES (GENDER / CONCERN / STAKEHOLDER, vertical-wise) ───────
    TS.add_template_slides(prs, a, b, label_a=fy1_label, label_b=fy2_label)

    # ── SLIDE 12: INSIGHTS ────────────────────────────────────────────────────
    s12 = CM.add_slide(prs)
    CM.add_logo(s12)
    CM.add_section_header(s12, "YEARLY KEY INSIGHTS", Inches(0.15), Inches(0.75),
                          Inches(13.0), Inches(0.45), bg_color=C.PRIMARY)
    insights = _generate_insights(a, b, v_labels, fy1_label, fy2_label)
    y = Inches(1.40)
    for i, line in enumerate(insights[:8], 1):
        CM.add_text(s12, f"{i}. {line}", Inches(0.8), y, Inches(11.7), Inches(0.45),
                    font_size=Pt(12), bold=False, color=C.BLACK, alignment=PP_ALIGN.LEFT)
        y += Inches(0.55)

    # ── SLIDE 13: THANK YOU ───────────────────────────────────────────────────
    s13 = CM.add_slide(prs)
    CM.add_text(s13, "Thank You", Inches(0.14), Inches(2.60), Inches(13.09), Inches(0.63),
                font_size=Pt(36), bold=True, color=C.PRIMARY, alignment=PP_ALIGN.CENTER)
    CM.add_text(s13, f"Wellness Annual Report -- {fy1_label} vs {fy2_label}",
                Inches(0.14), Inches(3.35), Inches(13.09), Inches(0.40),
                font_size=Pt(16), bold=False, color=C.DARK_GRAY, alignment=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _pct(old, new):
    if old == 0:
        return "N/A"
    return f"{((new - old) / old) * 100:+.1f}%"


def _generate_insights(a, b, v_labels, fy1, fy2):
    vk = C.VERT_KEYS
    lines = []
    delta = b["grand"] - a["grand"]
    pct = f" ({delta / a['grand'] * 100:+.1f}%)" if a["grand"] else ""
    lines.append(f"Grand total cases changed from {a['grand']} to {b['grand']} ({delta:+d}{pct}).")
    for key, name in [("new", "New cases"), ("followup", "Follow-up cases")]:
        d = b[key] - a[key]
        lines.append(f"{name}: {a[key]} to {b[key]} ({d:+d}).")
    n_verts = len(vk)
    largest_i = max(range(n_verts), key=lambda i: b["vertical"][vk[i]]["total"]
                    - a["vertical"][vk[i]]["total"])
    lines.append(f"Largest vertical increase: {v_labels[largest_i]} "
                 f"({a['vertical'][vk[largest_i]]['total']} "
                 f"to {b['vertical'][vk[largest_i]]['total']}).")
    smallest_i = min(range(n_verts), key=lambda i: b["vertical"][vk[i]]["total"]
                     - a["vertical"][vk[i]]["total"])
    lines.append(f"Largest vertical decrease: {v_labels[smallest_i]} "
                 f"({a['vertical'][vk[smallest_i]]['total']} "
                 f"to {b['vertical'][vk[smallest_i]]['total']}).")
    return lines
