"""Tests for PPTX deck generation (Phase 6)."""

import io

import pytest
from django.core.files.uploadedfile import SimpleUploadedFile
from pptx import Presentation

from tests.builders import fields
from wellness.models import Period, User
from wellness.services.persistence import save_import
from wellness.services.reports import ppt


@pytest.fixture
def admin_user(db):
    return User.objects.create_user(username="a1", password="x", role=User.Role.ADMIN)


def _save(admin_user, report_type, start, end, new, fu, secondary=None):
    from tests.test_api import workbook_bytes
    data = workbook_bytes(
        report_type=report_type, start=start, end=end,
        new_rows=new, fu_rows=fu, secondary=secondary,
    )
    from wellness.services import parsing as P
    rep = P.parse_excel(data)
    return save_import(rep, admin_user, f"{start}.xlsx", data).period


def _row(total=0, m=0, f=0, o=0, **kw):
    d = fields()
    d.update({"total_cases": total, "gender_male": m, "gender_female": f, "gender_other": o})
    d.update(kw)
    return d


@pytest.fixture
def weekly_pair(db, admin_user):
    a = _save(
        admin_user, "weekly", "15th July", "21st July 2026",
        {"WLN Ctr": _row(4, 2, 2, mode_online=4, concern_anxiety=4, stake_ug=4, referral_self=4),
         "Team A": _row(), "Your Dost": _row(), "Myndwell": _row()},
        {"WLN Ctr": _row(33, 20, 13), "Team A": _row(), "Your Dost": _row(), "Myndwell": _row()},
    )
    b = _save(
        admin_user, "weekly", "22nd July", "28th July 2026",
        {"WLN Ctr": _row(9, 5, 4), "Team A": _row(), "Your Dost": _row(3, 2, 1), "Myndwell": _row(7, 4, 3)},
        {"WLN Ctr": _row(32, 18, 14), "Team A": _row(), "Your Dost": _row(38, 20, 18), "Myndwell": _row(14, 8, 6)},
    )
    return a, b


@pytest.fixture
def monthly_period(db, admin_user):
    return _save(
        admin_user, "monthly", "01st July", "30th July 2026",
        {"WLN Ctr": _row(27, 15, 12), "Team A": _row(), "Your Dost": _row(22, 12, 10), "Myndwell": _row(23, 10, 13)},
        {"WLN Ctr": _row(101, 55, 46), "Team A": _row(), "Your Dost": _row(170, 90, 80), "Myndwell": _row(41, 20, 21)},
    )


class TestWeekly:
    def test_builds_valid_pptx(self, weekly_pair):
        a, b = weekly_pair
        buf = ppt.build_weekly(a, b)
        assert buf[:2] == b"PK"  # valid zip/pptx
        prs = Presentation(io.BytesIO(buf))
        slides = list(prs.slides)
        assert len(slides) == 10  # programmatic engine: 10 slides

    def test_weekly_header_contains_both_dates(self, weekly_pair):
        a, b = weekly_pair
        buf = ppt.build_weekly(a, b)
        prs = Presentation(io.BytesIO(buf))
        # Slide 1 header must contain both period date strings
        all_text = " ".join(
            sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame
        )
        assert "15th" in all_text and "21st" in all_text
        assert "22nd" in all_text and "28th" in all_text

    def test_weekly_footer_present_on_every_data_slide(self, weekly_pair):
        a, b = weekly_pair
        buf = ppt.build_weekly(a, b)
        prs = Presentation(io.BytesIO(buf))
        # Slides 1–10 must have "NEW CASES" in their text (footer strip)
        for i, slide in enumerate(list(prs.slides)[:10]):
            texts = " ".join(
                sh.text_frame.text for sh in slide.shapes if sh.has_text_frame
            )
            assert "NEW CASES" in texts, f"Slide {i + 1} missing NEW CASES footer"

    def test_weekly_month_cross(self, db, admin_user, weekly_pair):
        a, b = weekly_pair
        a.period_start = a.period_start.replace(month=7, day=29)
        a.period_end = a.period_end.replace(month=8, day=4)
        a.save()
        buf = ppt.build_weekly(a, b)
        prs = Presentation(io.BytesIO(buf))
        t = " ".join(sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame)
        assert "29th" in t and "4th" in t  # cross-month dates appear


class TestMonthly:
    def test_builds_valid_pptx(self, monthly_period):
        p = monthly_period
        buf = ppt.build_monthly(p)
        assert buf[:2] == b"PK"
        prs = Presentation(io.BytesIO(buf))
        slides = list(prs.slides)
        assert len(slides) == 13  # programmatic engine: 13 slides

    def test_monthly_header_contains_month(self, monthly_period):
        p = monthly_period
        buf = ppt.build_monthly(p)
        prs = Presentation(io.BytesIO(buf))
        t = " ".join(sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame)
        assert "JULY 2026" in t

    def test_monthly_footer_present_on_every_data_slide(self, monthly_period):
        p = monthly_period
        buf = ppt.build_monthly(p)
        prs = Presentation(io.BytesIO(buf))
        # Slides 1–12 must have "NEW CASES" in footer
        for i, slide in enumerate(list(prs.slides)[:12]):
            texts = " ".join(
                sh.text_frame.text for sh in slide.shapes if sh.has_text_frame
            )
            assert "NEW CASES" in texts, f"Monthly slide {i + 1} missing NEW CASES footer"

    def test_monthly_all_slide_titles_present(self, monthly_period):
        """Every major section slide must be identifiable by title text."""
        p = monthly_period
        buf = ppt.build_monthly(p)
        prs = Presentation(io.BytesIO(buf))
        all_text = " ".join(
            sh.text_frame.text
            for slide in prs.slides
            for sh in slide.shapes
            if sh.has_text_frame
        )
        for expected in ["MONTHLY WELLNESS DATA", "VERTICAL DISTRIBUTION",
                         "GENDER DISTRIBUTION", "MODE OF SESSION",
                         "REFERRAL TYPE", "RANGE OF CONCERN", "STAKEHOLDER",
                         "GRAND TOTAL", "AI DATA INSIGHTS"]:
            assert expected in all_text, f"Missing slide title: {expected}"

    def test_monthly_different_month(self, db, admin_user, monthly_period):
        p = monthly_period
        p.period_end = p.period_end.replace(month=8, day=31)
        p.period_start = p.period_start.replace(month=8, day=1)
        p.save()
        buf = ppt.build_monthly(p)
        prs = Presentation(io.BytesIO(buf))
        t = " ".join(sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame)
        assert "AUGUST 2026" in t


class TestYearly:
    def test_builds_twelve_slide_wlc_team_a_deck(self, monthly_period):
        buf = ppt.build_yearly([monthly_period], [monthly_period])
        assert buf[:2] == b"PK"
        prs = Presentation(io.BytesIO(buf))
        assert len(prs.slides) == 12
        all_text = " ".join(
            sh.text_frame.text
            for slide in prs.slides
            for sh in slide.shapes
            if sh.has_text_frame
        )
        assert "Team A" in all_text
        assert "WC" in all_text
        assert "STAKEHOLDER" in all_text
        assert "RANGE OF CONCERN ADDRESSED" in all_text



class TestExports:
    def test_excel(self, monthly_period):
        from wellness.services.reports import exports
        data = exports.build_excel(monthly_period)
        assert data[:2] == b"PK"
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(data))
        assert "Summary" in wb.sheetnames
        ws = wb["Summary"]
        assert ws["B5"].value == "WLN Ctr"
        assert ws["B6"].value == "Team A"
        assert ws["C15"].value == 384  # grand total after four header rows and both blocks
        assert ws["A17"].value == "Unrecognised"  # secondary block is on Summary

    def test_csv(self, monthly_period):
        from wellness.services.reports import exports
        text = exports.build_csv(monthly_period).decode("utf-8-sig")
        assert "new,WC,27,15,12" in text
        assert "new,TA,0,0,0" in text
        assert "followup" in text

    def test_pdf(self, monthly_period):
        from wellness.services.reports import exports
        data = exports.build_pdf(monthly_period)
        assert data[:4] == b"%PDF"

    def test_build_dispatch(self, monthly_period):
        from wellness.services.reports import exports
        fn, data, ct = exports.build("xlsx", monthly_period)
        assert fn.endswith(".xlsx") and data[:2] == b"PK"


class TestReportEndpoint:
    def _client_period(self, db, admin_user):
        from rest_framework.test import APIClient
        from django.core.files.uploadedfile import SimpleUploadedFile
        import tests.test_api as ta
        c = APIClient()
        c.force_authenticate(user=admin_user)
        data = ta.workbook_bytes(
            report_type="monthly", start="01st July", end="30th July 2026",
            new_rows={t: ta.good_row() for t in ("WLN Ctr", "Team A", "Your Dost", "Myndwell")},
            fu_rows={t: ta.good_row() for t in ("WLN Ctr", "Team A", "Your Dost", "Myndwell")},
        )
        f = SimpleUploadedFile("m.xlsx", data)
        p = c.post("/api/imports/preview", {"file": f}, format="multipart")
        confirm = c.post("/api/imports/confirm", {"preview_id": p.data["preview_id"]}, format="json")
        return c, confirm.data["period_id"]

    def test_ppt_monthly(self, db, admin_user):
        c, pid = self._client_period(db, admin_user)
        r = c.post("/api/reports/generate", {"period_id": pid, "format": "ppt"}, format="json")
        assert r.status_code == 200
        assert r["Content-Disposition"]
        assert b"".join(r.streaming_content)[:2] == b"PK"

    def test_ppt_weekly_auto_previous(self, db, admin_user):
        from rest_framework.test import APIClient
        from django.core.files.uploadedfile import SimpleUploadedFile
        import tests.test_api as ta
        c = APIClient()
        c.force_authenticate(user=admin_user)
        for start, end in (("15th July", "21st July 2026"), ("22nd July", "28th July 2026")):
            data = ta.workbook_bytes(report_type="weekly", start=start, end=end,
                                     new_rows={t: ta.good_row() for t in ("WLN Ctr", "Team A", "Your Dost", "Myndwell")},
                                     fu_rows={t: ta.good_row() for t in ("WLN Ctr", "Team A", "Your Dost", "Myndwell")})
            f = SimpleUploadedFile("w.xlsx", data)
            p = c.post("/api/imports/preview", {"file": f}, format="multipart")
            c.post("/api/imports/confirm", {"preview_id": p.data["preview_id"]}, format="json")
        cur = Period.objects.get(period_start="2026-07-22")
        r = c.post("/api/reports/generate", {"period_id": cur.id, "format": "ppt"}, format="json")
        assert r.status_code == 200
        prs = Presentation(io.BytesIO(b"".join(r.streaming_content)))
        t = " ".join(sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame)
        # Both period dates must appear somewhere in slide 1 text
        assert "15th" in t and "22nd" in t

    def test_excel_and_csv_via_api(self, db, admin_user):
        c, pid = self._client_period(db, admin_user)
        assert c.post("/api/reports/generate", {"period_id": pid, "format": "xlsx"}, format="json").status_code == 200
        assert c.post("/api/reports/generate", {"period_id": pid, "format": "csv"}, format="json").status_code == 200
        assert c.post("/api/reports/generate", {"period_id": pid, "format": "pdf"}, format="json").status_code == 200

    def test_missing_period(self, db, admin_user):
        from rest_framework.test import APIClient
        c = APIClient()
        c.force_authenticate(user=admin_user)
        assert c.post("/api/reports/generate", {"period_id": 999}, format="json").status_code == 400


class TestReportCenter:
    def test_single_reports_are_real_office_files(self, weekly_pair):
        from wellness.services.reports import report_center

        period, _ = weekly_pair
        for fmt, extension, content_type in (
            ("ppt", ".pptx", "presentationml"),
            ("xlsx", ".xlsx", "spreadsheetml"),
        ):
            filename, data, actual_type, source_ids = report_center.build_single(
                "weekly", fmt, period_id=period.id,
            )
            assert filename.endswith(extension)
            assert data[:2] == b"PK"
            assert content_type in actual_type
            assert source_ids == [period.id]

    def test_week_compare_generates_pptx(self, weekly_pair):
        from wellness.services.reports import report_center

        first, second = weekly_pair
        filename, data, content_type = report_center.build_compare(
            "week", "ppt", from_id=first.id, to_id=second.id,
        )
        assert filename.endswith(".pptx")
        assert data[:2] == b"PK"
        assert content_type.endswith("presentationml.presentation")
