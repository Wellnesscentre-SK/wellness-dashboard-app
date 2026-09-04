"""
AI Insights & Improvement Recommendations PPT slides.

Adds professionally designed slides to existing PPT reports matching
the existing Wellness Centre design system (blue headers, Calibri, 16:9).
"""

from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
try:
    from ppt_generator import config as C
    from ppt_generator import components as CM
except ImportError:
    import config as C
    import components as CM


PRIORITY_COLORS = {
    "HIGH": C.DANGER,
    "MEDIUM": C.WARNING,
    "LOW": C.SUCCESS,
}

CATEGORY_ICONS = {
    "Performance Improvement": "\u25B2",
    "Client Engagement": "\u2665",
    "Operational Improvement": "\u2699",
    "Team Development": "\u2605",
    "Wellness Centre Development": "\u271A",
    "Reporting & Data Quality": "\u2630",
    "Future Opportunities": "\u2728",
}


def _add_ai_header(slide, title, subtitle=""):
    """Add the standard blue header bar with title."""
    CM.add_logo(slide)
    CM.add_section_header(
        slide, title,
        Inches(0.15), Inches(0.75), Inches(13.0), Inches(0.45),
        bg_color=C.PRIMARY,
    )
    if subtitle:
        CM.add_text(
            slide, subtitle,
            Inches(0.15), Inches(1.25), Inches(13.0), Inches(0.30),
            font_size=C.SMALL_FONT, bold=False, color=C.GRAY,
            alignment=PP_ALIGN.CENTER,
        )


def _add_priority_badge(slide, left, top, priority):
    """Add a small colored priority indicator."""
    color = PRIORITY_COLORS.get(priority, C.GRAY)
    shape = slide.shapes.add_shape(1, left, top, Inches(0.8), Inches(0.25))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = priority
    run.font.size = Pt(8)
    run.font.bold = True
    run.font.color.rgb = C.WHITE
    run.font.name = C.FONT_NAME


def _add_footer(slide):
    CM.add_footer(slide)


# ---------------------------------------------------------------------------
# Slide 1: AI Insights & Improvement Recommendations (Summary)
# ---------------------------------------------------------------------------

def add_ai_summary_slide(prs, ai_result, period_label=""):
    """Add the main AI insights summary slide."""
    slide = CM.add_slide(prs)
    _add_ai_header(
        slide,
        "AI INSIGHTS & IMPROVEMENT RECOMMENDATIONS",
        f"AI-generated insights based on {period_label} performance data" if period_label else "",
    )

    summary = ai_result.get("summary", "")
    kpi = ai_result.get("kpi", {})

    # Summary text box
    if summary:
        CM.add_text(
            slide, f"AI Summary: {summary}",
            Inches(0.4), Inches(1.65), Inches(12.5), Inches(0.65),
            font_size=Pt(11), bold=False, color=C.DARK_GRAY,
            alignment=PP_ALIGN.LEFT,
        )

    # KPI cards row
    y_kpi = Inches(2.50)
    cards = [
        ("AI INSIGHTS", kpi.get("total_suggestions", 0), C.PRIMARY),
        ("HIGH PRIORITY", kpi.get("high", 0), C.DANGER),
        ("MEDIUM PRIORITY", kpi.get("medium", 0), C.WARNING),
        ("LOW PRIORITY", kpi.get("low", 0), C.SUCCESS),
        ("OPPORTUNITIES", kpi.get("opportunities", 0), C.INFO),
    ]
    CM.add_kpi_cards(slide, cards, top=y_kpi)

    # Top 5 recommendations
    y = Inches(3.60)
    CM.add_text(
        slide, "TOP RECOMMENDATIONS",
        Inches(0.4), y, Inches(12.5), Inches(0.35),
        font_size=C.HEADING_FONT, bold=True, color=C.PRIMARY,
        alignment=PP_ALIGN.LEFT,
    )
    y += Inches(0.40)

    recs = ai_result.get("suggestions", [])[:5]
    for i, rec in enumerate(recs, 1):
        priority = rec.get("priority", "MEDIUM")
        title = rec.get("title", "")
        category = rec.get("category_label", "")
        evidence = rec.get("evidence", "")

        # Recommendation number + title
        CM.add_text(
            slide, f"{i}.",
            Inches(0.4), y, Inches(0.35), Inches(0.30),
            font_size=Pt(12), bold=True, color=C.PRIMARY,
            alignment=PP_ALIGN.RIGHT,
        )
        CM.add_text(
            slide, title,
            Inches(0.8), y, Inches(8.5), Inches(0.30),
            font_size=Pt(12), bold=True, color=C.BLACK,
            alignment=PP_ALIGN.LEFT,
        )
        _add_priority_badge(slide, Inches(9.5), y + Inches(0.03), priority)
        CM.add_text(
            slide, category,
            Inches(10.5), y, Inches(2.5), Inches(0.25),
            font_size=Pt(9), bold=False, color=C.GRAY,
            alignment=PP_ALIGN.LEFT,
        )

        # Evidence line
        if evidence:
            CM.add_text(
                slide, f"Evidence: {evidence}",
                Inches(0.8), y + Inches(0.30), Inches(12.2), Inches(0.25),
                font_size=Pt(9), bold=False, color=C.GRAY,
                alignment=PP_ALIGN.LEFT,
            )
        y += Inches(0.60)
        if y > Inches(6.8):
            break

    _add_footer(slide)
    return slide


# ---------------------------------------------------------------------------
# Slide 2: Detailed Recommendations by Category
# ---------------------------------------------------------------------------

def add_ai_detail_slide(prs, suggestions, category_filter=None, slide_title="AI Recommendations"):
    """Add a detailed recommendations slide, optionally filtered by category."""
    slide = CM.add_slide(prs)
    _add_ai_header(slide, slide_title)

    filtered = suggestions
    if category_filter:
        filtered = [s for s in suggestions if s.get("category") == category_filter]

    y = Inches(1.65)
    for i, rec in enumerate(filtered[:6], 1):
        priority = rec.get("priority", "MEDIUM")
        title = rec.get("title", "")
        why = rec.get("why", "")
        action = rec.get("action", "")
        benefit = rec.get("benefit", "")

        # Title + priority
        CM.add_text(
            slide, f"{i}. {title}",
            Inches(0.4), y, Inches(10.5), Inches(0.28),
            font_size=Pt(11), bold=True, color=C.BLACK,
            alignment=PP_ALIGN.LEFT,
        )
        _add_priority_badge(slide, Inches(11.2), y + Inches(0.02), priority)

        # Why
        if why:
            CM.add_text(
                slide, f"Why: {why}",
                Inches(0.6), y + Inches(0.28), Inches(12.4), Inches(0.22),
                font_size=Pt(9), bold=False, color=C.DARK_GRAY,
                alignment=PP_ALIGN.LEFT,
            )
        # Action
        if action:
            CM.add_text(
                slide, f"Action: {action}",
                Inches(0.6), y + Inches(0.50), Inches(12.4), Inches(0.22),
                font_size=Pt(9), bold=False, color=C.GRAY,
                alignment=PP_ALIGN.LEFT,
            )
        # Benefit
        if benefit:
            CM.add_text(
                slide, f"Benefit: {benefit}",
                Inches(0.6), y + Inches(0.72), Inches(12.4), Inches(0.22),
                font_size=Pt(9), bold=False, color=C.SUCCESS,
                alignment=PP_ALIGN.LEFT,
            )

        y += Inches(0.98)
        if y > Inches(6.8):
            break

    _add_footer(slide)
    return slide


# ---------------------------------------------------------------------------
# Slide 3: Improvement Roadmap
# ---------------------------------------------------------------------------

def add_ai_roadmap_slide(prs, roadmap, period_label=""):
    """Add an improvement roadmap slide with 4 time horizons."""
    slide = CM.add_slide(prs)
    _add_ai_header(
        slide,
        "AI IMPROVEMENT ROADMAP",
        f"Strategic improvement timeline for {period_label}" if period_label else "",
    )

    phases = [
        ("immediate", "IMMEDIATE\n0-7 Days", C.DANGER),
        ("short_term", "SHORT TERM\n1-4 Weeks", C.WARNING),
        ("medium_term", "MEDIUM TERM\n1-3 Months", C.INFO),
        ("long_term", "LONG TERM\n3-12 Months", C.SUCCESS),
    ]

    col_w = Inches(3.0)
    gap = Inches(0.2)
    start_x = Inches(0.4)

    for idx, (key, label, color) in enumerate(phases):
        x = start_x + idx * (col_w + gap)
        y_top = Inches(1.70)

        # Phase header
        shape = slide.shapes.add_shape(1, x, y_top, col_w, Inches(0.55))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = C.WHITE
        run.font.name = C.FONT_NAME

        # Items
        items = roadmap.get(key, {}).get("items", []) if isinstance(roadmap.get(key), dict) else []
        y_item = y_top + Inches(0.65)
        for item in items[:4]:
            title = item.get("title", "")
            priority = item.get("priority", "")
            if title:
                CM.add_text(
                    slide, f"\u2022 {title}",
                    x + Inches(0.1), y_item, col_w - Inches(0.2), Inches(0.35),
                    font_size=Pt(9), bold=False, color=C.BLACK,
                    alignment=PP_ALIGN.LEFT,
                )
                y_item += Inches(0.35)
        if not items:
            CM.add_text(
                slide, "No items in this phase.",
                x + Inches(0.1), y_item, col_w - Inches(0.2), Inches(0.3),
                font_size=Pt(9), bold=False, color=C.GRAY,
                alignment=PP_ALIGN.CENTER,
            )

    _add_footer(slide)
    return slide


# ---------------------------------------------------------------------------
# Slide 4: Next Actions
# ---------------------------------------------------------------------------

def add_ai_actions_slide(prs, suggestions, period_label=""):
    """Add a recommended next actions slide."""
    slide = CM.add_slide(prs)
    _add_ai_header(
        slide,
        "RECOMMENDED NEXT ACTIONS",
        f"Priority actions for {period_label}" if period_label else "",
    )

    high = [s for s in suggestions if s.get("priority") == "HIGH"]
    medium = [s for s in suggestions if s.get("priority") == "MEDIUM"]

    y = Inches(1.70)
    action_num = 1

    if high:
        CM.add_text(
            slide, "HIGH PRIORITY ACTIONS",
            Inches(0.4), y, Inches(12.5), Inches(0.30),
            font_size=C.SECTION_FONT, bold=True, color=C.DANGER,
            alignment=PP_ALIGN.LEFT,
        )
        y += Inches(0.40)
        for rec in high[:3]:
            CM.add_text(
                slide, f"{action_num}.",
                Inches(0.5), y, Inches(0.35), Inches(0.28),
                font_size=Pt(12), bold=True, color=C.DANGER,
                alignment=PP_ALIGN.RIGHT,
            )
            CM.add_text(
                slide, rec.get("action", rec.get("title", "")),
                Inches(0.9), y, Inches(12.0), Inches(0.28),
                font_size=Pt(11), bold=False, color=C.BLACK,
                alignment=PP_ALIGN.LEFT,
            )
            y += Inches(0.35)
            action_num += 1
        y += Inches(0.15)

    if medium:
        CM.add_text(
            slide, "MEDIUM PRIORITY ACTIONS",
            Inches(0.4), y, Inches(12.5), Inches(0.30),
            font_size=C.SECTION_FONT, bold=True, color=C.WARNING,
            alignment=PP_ALIGN.LEFT,
        )
        y += Inches(0.40)
        for rec in medium[:3]:
            CM.add_text(
                slide, f"{action_num}.",
                Inches(0.5), y, Inches(0.35), Inches(0.28),
                font_size=Pt(12), bold=True, color=C.WARNING,
                alignment=PP_ALIGN.RIGHT,
            )
            CM.add_text(
                slide, rec.get("action", rec.get("title", "")),
                Inches(0.9), y, Inches(12.0), Inches(0.28),
                font_size=Pt(11), bold=False, color=C.BLACK,
                alignment=PP_ALIGN.LEFT,
            )
            y += Inches(0.35)
            action_num += 1

    _add_footer(slide)
    return slide


# ---------------------------------------------------------------------------
# Public API: add all AI slides
# ---------------------------------------------------------------------------

def add_ai_slides(prs, ai_result, period_label="", max_slides=4):
    """Add all AI insight slides to a presentation.

    Parameters
    ----------
    prs : Presentation
        The presentation to append slides to.
    ai_result : dict
        Output from generate_*_suggestions() — must contain 'summary',
        'suggestions', 'kpi', and 'roadmap'.
    period_label : str
        Human-readable label for the period being analysed.
    max_slides : int
        Maximum number of AI slides to add (1-4).
    """
    suggestions = ai_result.get("suggestions", [])
    roadmap = ai_result.get("roadmap", {})

    if not suggestions:
        return

    add_ai_summary_slide(prs, ai_result, period_label)

    if max_slides >= 2 and suggestions:
        add_ai_detail_slide(prs, suggestions, slide_title="AI Recommendations \u2014 Performance & Operations")

    if max_slides >= 3 and roadmap:
        add_ai_roadmap_slide(prs, roadmap, period_label)

    if max_slides >= 4:
        add_ai_actions_slide(prs, suggestions, period_label)
