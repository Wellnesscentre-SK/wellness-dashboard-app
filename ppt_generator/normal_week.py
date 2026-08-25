"""
Normal Weekly Report — single week snapshot.
Excel-style layout: KPI cards, verticals table+pie, demographics, concern+pie, stakeholder+pie, bar chart.
"""

import io
from pptx.util import Inches, Pt
import config as C
import components as CM
import template_slides as TS


def build(period: dict, insights: list = None) -> bytes:
    prs = CM.create_presentation()
    label = period.get("label", "Week")
    grand = period.get("grand", 0)
    new_n = period.get("new", 0)
    fu_n = period.get("followup", 0)
    p = {**period, "vertical": C.combine_verticals(period.get("vertical", {}))}
    vk = C.VERT_KEYS

    # ── SLIDE 1: COVER ────────────────────────────────────────────────────────
    CM.add_cover(prs, "Weekly Wellness Data Report", label)

    # ── SLIDE 2: KPI CARDS + VERTICALS TABLE + PIE ────────────────────────────
    s2 = CM.add_slide(prs)
    CM.add_logo(s2)
    CM.add_kpi_cards(s2, [
        ("Total Cases", grand, C.TOTAL_CLR),
        ("New Cases", new_n, C.NEW_CASE_CLR),
        ("Follow-up Cases", fu_n, C.FOLLOWUP_CLR),
    ])

    v_labels = [C.VERTICALS[v] for v in vk]
    v_totals = [p["vertical"][v]["total"] for v in vk]
    v_new = [p["vertical"][v]["new"] for v in vk]
    v_fu = [p["vertical"][v]["followup"] for v in vk]

    vtable = [["Vertical", "Total Cases", "% Share"]]
    for i, v in enumerate(vk):
        vtable.append([C.VERTICALS[v], str(v_totals[i]),
                        f"{v_totals[i]/grand*100:.0f}%" if grand else "0%"])

    add_table_slide = CM.add_table(s2, Inches(0.25), Inches(1.85),
                                    Inches(5.5), Inches(1.5), vtable)
    CM.add_pie_chart(s2, Inches(6.5), Inches(1.70), Inches(6.5), Inches(4.5),
                     v_labels, v_totals)

    # ── SLIDE 3: VERTICALS NEW vs FOLLOW-UP TABLE + BAR CHART ─────────────────
    s3 = CM.add_slide(prs)
    CM.add_logo(s3)
    CM.add_section_header(s3, "NEW vs FOLLOW-UP BY VERTICAL", Inches(0.15), Inches(0.75),
                          Inches(13.0), Inches(0.45), bg_color=C.PRIMARY)
    nf_table = [["Vertical", "New Cases", "Follow-up", "New %", "F/up %"]]
    for i, v in enumerate(vk):
        t = v_totals[i] or 1
        nf_table.append([C.VERTICALS[v], str(v_new[i]), str(v_fu[i]),
                          f"{v_new[i]/t*100:.0f}%", f"{v_fu[i]/t*100:.0f}%"])
    CM.add_table(s3, Inches(0.25), Inches(1.35), Inches(6.0), Inches(2.0), nf_table)
    CM.add_column_chart(s3, Inches(6.8), Inches(1.20), Inches(6.2), Inches(5.5),
                        v_labels, [("New", v_new), ("Follow-up", v_fu)],
                        legend_position="bottom",
                        series_colors=[C.NEW_CASE_CLR, C.FOLLOWUP_CLR])

    # ── SLIDE 4: DEMOGRAPHICS & MODE ──────────────────────────────────────────
    s4 = CM.add_slide(prs)
    CM.add_logo(s4)
    CM.add_section_header(s4, "DEMOGRAPHICS & MODE OF SESSION", Inches(0.15), Inches(0.75),
                          Inches(13.0), Inches(0.45), bg_color=C.PRIMARY)

    g_table = [["Gender", "Count", "% Share"]]
    for g in C.GENDER_LABELS:
        gval = p["gender"].get(g, 0)
        g_table.append([g, str(gval), f"{gval/grand*100:.0f}%" if grand else "0%"])

    m_table = [["Mode", "Count", "% Share"]]
    for m in C.MODE_LABELS:
        mval = p["mode"].get(m, 0)
        m_table.append([m, str(mval), f"{mval/grand*100:.0f}%" if grand else "0%"])

    CM.add_table(s4, Inches(0.25), Inches(1.35), Inches(4.0), Inches(2.2), g_table)
    CM.add_table(s4, Inches(4.50), Inches(1.35), Inches(4.0), Inches(2.2), m_table)

    r_table = [["Referral", "Count", "% Share"]]
    for r in C.REFERRAL_LABELS:
        rval = p["referral"].get(r, 0)
        r_table.append([r, str(rval), f"{rval/grand*100:.0f}%" if grand else "0%"])
    CM.add_table(s4, Inches(8.75), Inches(1.35), Inches(4.3), Inches(3.0), r_table)

    # ── SLIDE 5: CONCERN TABLE + PIE ──────────────────────────────────────────
    s5 = CM.add_slide(prs)
    CM.add_logo(s5)
    CM.add_section_header(s5, "RANGE OF CONCERN ADDRESSED", Inches(0.15), Inches(0.75),
                          Inches(13.0), Inches(0.45), bg_color=C.PRIMARY)
    c_table = [["Concern", "Count", "% Share"]]
    c_labels = C.CONCERN_LABELS
    c_values = [p["concern"].get(c, 0) for c in c_labels]
    for i, c in enumerate(c_labels):
        c_table.append([c, str(c_values[i]),
                        f"{c_values[i]/grand*100:.0f}%" if grand else "0%"])
    CM.add_table(s5, Inches(0.25), Inches(1.35), Inches(6.0), Inches(4.5), c_table)
    CM.add_pie_chart(s5, Inches(6.8), Inches(1.35), Inches(6.2), Inches(5.0),
                     c_labels, c_values)

    # ── SLIDE 6: STAKEHOLDER TABLE + PIE ──────────────────────────────────────
    s6 = CM.add_slide(prs)
    CM.add_logo(s6)
    CM.add_section_header(s6, "STAKEHOLDER", Inches(0.15), Inches(0.75),
                          Inches(13.0), Inches(0.45), bg_color=C.PRIMARY)
    st_labels = C.STAKEHOLDER_LABELS
    st_values = [p["stakeholder"].get(s, 0) for s in st_labels]
    st_table = [["Stakeholder", "Count", "% Share"]]
    for i, s in enumerate(st_labels):
        st_table.append([s, str(st_values[i]),
                         f"{st_values[i]/grand*100:.0f}%" if grand else "0%"])
    CM.add_table(s6, Inches(0.25), Inches(1.35), Inches(6.0), Inches(4.5), st_table)
    CM.add_pie_chart(s6, Inches(6.8), Inches(1.35), Inches(6.2), Inches(5.0),
                     st_labels, st_values)

    # ── SLIDES 8-10: TEMPLATE (GENDER / CONCERN / STAKEHOLDER, vertical-wise) ─
    TS.add_template_slides(prs, p)

    # ── SLIDE 11: AI INSIGHTS ─────────────────────────────────────────────────
    if insights is None:
        import ai_analyzer
        insights = ai_analyzer.analyze(p)
    CM.add_insights_slide(prs, "AI DATA INSIGHTS", insights, label)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
