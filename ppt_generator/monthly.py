"""
Monthly Wellness Comparison Report — Excel-style layout.
KPI comparison, verticals, demographics, concern, stakeholder with tables+pies+bars.
"""

import io
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import config as C
import components as CM
import template_slides as TS


def build(period_a: dict, period_b: dict,
          key_insights: list = None, proposed_points: list = None) -> bytes:
    prs = CM.create_presentation()
    la = period_a.get("label", "Period A")
    lb = period_b.get("label", "Period B")
    va = {**period_a, "vertical": C.combine_verticals(period_a.get("vertical", {}))}
    vb = {**period_b, "vertical": C.combine_verticals(period_b.get("vertical", {}))}
    vk = C.VERT_KEYS

    # ── SLIDE 1: COVER ────────────────────────────────────────────────────────
    CM.add_cover(prs, "Monthly Wellness Data", f"Comparative Analysis: {la} vs. {lb}")

    # ── SLIDE 2: KPI COMPARISON + SUMMARY TABLES ──────────────────────────────
    s2 = CM.add_slide(prs)
    CM.add_logo(s2)
    CM.add_kpi_cards(s2, [
        (f"Total ({la})", va["grand"], C.TOTAL_CLR),
        (f"Total ({lb})", vb["grand"], C.NEW_CASE_CLR),
        ("Difference", vb["grand"] - va["grand"], C.FOLLOWUP_CLR),
        ("Top Vertical", "Your Dost", C.TOP_VERT_CLR),
    ])
    CM.add_table(s2, Inches(0.45), Inches(1.85), Inches(5.7), Inches(2.0), [
        ["Metric", la, lb, "Change"],
        ["New Cases", va["new"], vb["new"], vb["new"] - va["new"]],
        ["Follow-up Cases", va["followup"], vb["followup"], vb["followup"] - va["followup"]],
        ["Grand Total", va["grand"], vb["grand"], vb["grand"] - va["grand"]],
    ])
    CM.add_table(s2, Inches(6.50), Inches(1.85), Inches(6.3), Inches(2.0),
                 [["Gender", la, lb]] + [
        [g, va["gender"].get(g, 0), vb["gender"].get(g, 0)] for g in C.GENDER_LABELS])
    CM.add_table(s2, Inches(0.45), Inches(4.00), Inches(5.7), Inches(2.0),
                 [["Mode", la, lb]] + [
        [m, va["mode"].get(m, 0), vb["mode"].get(m, 0)] for m in C.MODE_LABELS])
    CM.add_table(s2, Inches(6.50), Inches(4.00), Inches(6.3), Inches(2.5),
                 [["Referral", la, lb]] + [
        [r, va["referral"].get(r, 0), vb["referral"].get(r, 0)] for r in C.REFERRAL_LABELS])

    # ── SLIDE 3: VERTICALS PIE PAIR ───────────────────────────────────────────
    v_labels = [C.VERTICALS[v] for v in vk]
    s3 = CM.add_slide(prs)
    CM.add_logo(s3)
    vals_a = [va["vertical"][v]["total"] for v in vk]
    vals_b = [vb["vertical"][v]["total"] for v in vk]
    CM.add_pie_pair(s3, v_labels, v_labels, vals_a, vals_b,
                    ["VERTICALS TOTAL", la], ["VERTICALS TOTAL", lb])

    # ── SLIDE 4: NEW vs FOLLOW-UP BAR (SIDE BY SIDE) ─────────────────────────
    s4 = CM.add_slide(prs)
    CM.add_logo(s4)
    CM.add_section_header(s4, "NEW vs FOLLOW-UP BY VERTICAL", Inches(0.15), Inches(0.75),
                          Inches(13.0), Inches(0.45), bg_color=C.PRIMARY)
    CM.add_text(s4, la, Inches(0.15), Inches(1.25), Inches(6.2), Inches(0.30),
                font_size=C.HEADING_FONT, bold=True, color=C.PRIMARY, alignment=PP_ALIGN.CENTER)
    CM.add_text(s4, lb, Inches(6.95), Inches(1.25), Inches(6.2), Inches(0.30),
                font_size=C.HEADING_FONT, bold=True, color=C.PRIMARY, alignment=PP_ALIGN.CENTER)
    CM.add_column_chart(s4, Inches(0.05), Inches(1.60), Inches(6.2), Inches(5.5),
                        v_labels,
                        [("New", [va["vertical"][v]["new"] for v in vk]),
                         ("Follow-up", [va["vertical"][v]["followup"] for v in vk])],
                        legend_position="bottom",
                        series_colors=[C.NEW_CASE_CLR, C.FOLLOWUP_CLR])
    CM.add_column_chart(s4, Inches(6.85), Inches(1.60), Inches(6.2), Inches(5.5),
                        v_labels,
                        [("New", [vb["vertical"][v]["new"] for v in vk]),
                         ("Follow-up", [vb["vertical"][v]["followup"] for v in vk])],
                        legend_position="bottom",
                        series_colors=[C.NEW_CASE_CLR, C.FOLLOWUP_CLR])

    # ── SLIDE 5: GENDER + MODE PIE PAIR ───────────────────────────────────────
    s5 = CM.add_slide(prs)
    CM.add_logo(s5)
    g_labels = C.GENDER_LABELS
    g_a = [va["gender"].get(g, 0) for g in g_labels]
    g_b = [vb["gender"].get(g, 0) for g in g_labels]
    CM.add_pie_pair(s5, g_labels, g_labels, g_a, g_b,
                    ["GENDER DISTRIBUTION", la], ["GENDER DISTRIBUTION", lb])

    s6 = CM.add_slide(prs)
    CM.add_logo(s6)
    m_labels = C.MODE_LABELS
    m_a = [va["mode"].get(m, 0) for m in m_labels]
    m_b = [vb["mode"].get(m, 0) for m in m_labels]
    CM.add_pie_pair(s6, m_labels, m_labels, m_a, m_b,
                    ["MODE OF SESSION", la], ["MODE OF SESSION", lb])

    # ── SLIDE 6: REFERRAL PIE PAIR ────────────────────────────────────────────
    s7 = CM.add_slide(prs)
    CM.add_logo(s7)
    r_labels = C.REFERRAL_LABELS
    r_a = [va["referral"].get(r, 0) for r in r_labels]
    r_b = [vb["referral"].get(r, 0) for r in r_labels]
    CM.add_pie_pair(s7, r_labels, r_labels, r_a, r_b,
                    ["REFERRAL TYPE", la], ["REFERRAL TYPE", lb])

    # ── SLIDE 7: CONCERN PIE PAIR ─────────────────────────────────────────────
    s8 = CM.add_slide(prs)
    CM.add_logo(s8)
    c_labels = C.CONCERN_LABELS
    c_a = [va["concern"].get(c, 0) for c in c_labels]
    c_b = [vb["concern"].get(c, 0) for c in c_labels]
    CM.add_pie_pair(s8, c_labels, c_labels, c_a, c_b,
                    ["RANGE OF CONCERN", la], ["RANGE OF CONCERN", lb])

    # ── SLIDE 8: STAKEHOLDER PIE PAIR ─────────────────────────────────────────
    s9 = CM.add_slide(prs)
    CM.add_logo(s9)
    st_labels = C.STAKEHOLDER_LABELS
    st_a = [va["stakeholder"].get(s, 0) for s in st_labels]
    st_b = [vb["stakeholder"].get(s, 0) for s in st_labels]
    CM.add_pie_pair(s9, st_labels, st_labels, st_a, st_b,
                    ["STAKEHOLDER", la], ["STAKEHOLDER", lb])

    # ── SLIDE 9: KEY INSIGHTS ─────────────────────────────────────────────────
    s10 = CM.add_slide(prs)
    CM.add_logo(s10)
    CM.add_section_header(s10, "KEY INSIGHTS", Inches(0.15), Inches(0.75),
                          Inches(13.0), Inches(0.45), bg_color=C.PRIMARY)
    if key_insights:
        y = Inches(1.40)
        for i, insight in enumerate(key_insights[:10], 1):
            CM.add_text(s10, f"{i}. {insight}", Inches(0.8), y,
                        Inches(11.7), Inches(0.45),
                        font_size=Pt(12), bold=False, color=C.BLACK,
                        alignment=PP_ALIGN.LEFT)
            y += Inches(0.50)

    # ── SLIDES 11-13: TEMPLATE (GENDER / CONCERN / STAKEHOLDER, vertical-wise) ─
    TS.add_template_slides(prs, va, vb, label_a=la, label_b=lb)

    CM.add_proposed_points(prs, proposed_points)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
