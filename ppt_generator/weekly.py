"""
Weekly Wellness Comparison Report — Excel-style layout.
KPI comparison cards, verticals table+pie (both periods), New vs Follow-up bar, concern, stakeholder.
"""

import io
from pptx.util import Inches, Pt
import config as C
import components as CM
import template_slides as TS


def build(period_a: dict, period_b: dict, proposed_points: list = None) -> bytes:
    prs = CM.create_presentation()
    la = period_a.get("label", "Period A")
    lb = period_b.get("label", "Period B")
    va = {**period_a, "vertical": C.combine_verticals(period_a.get("vertical", {}))}
    vb = {**period_b, "vertical": C.combine_verticals(period_b.get("vertical", {}))}
    vk = C.VERT_KEYS

    # ── SLIDE 1: COVER ────────────────────────────────────────────────────────
    CM.add_cover(prs, "Comparative Weekly Wellness Data", f"{la}  to  {lb}")

    # ── SLIDE 2: KPI COMPARISON + SUMMARY TABLE ───────────────────────────────
    s2 = CM.add_slide(prs)
    CM.add_logo(s2)
    CM.add_kpi_cards(s2, [
        (f"Total ({la})", va.get("grand", 0), C.TOTAL_CLR),
        (f"Total ({lb})", vb.get("grand", 0), C.NEW_CASE_CLR),
        ("Difference", vb.get("grand", 0) - va.get("grand", 0), C.FOLLOWUP_CLR),
        (f"New ({la})", va.get("new", 0), C.TOTAL_CLR),
        (f"New ({lb})", vb.get("new", 0), C.NEW_CASE_CLR),
    ])
    CM.add_table(s2, Inches(0.45), Inches(1.85), Inches(12.4), Inches(2.0), [
        ["Metric", la, lb, "Difference"],
        ["New Cases", va.get("new", 0), vb.get("new", 0), vb.get("new", 0) - va.get("new", 0)],
        ["Follow-up Cases", va.get("followup", 0), vb.get("followup", 0),
         vb.get("followup", 0) - va.get("followup", 0)],
        ["Grand Total", va.get("grand", 0), vb.get("grand", 0),
         vb.get("grand", 0) - va.get("grand", 0)],
    ])
    CM.add_table(s2, Inches(0.45), Inches(4.00), Inches(12.4), Inches(2.5), [
        ["Vertical", f"New {la}", f"New {lb}", f"Follow-up {la}", f"Follow-up {lb}"],
        *[[C.VERTICALS[v], va["vertical"][v]["new"], vb["vertical"][v]["new"],
           va["vertical"][v]["followup"], vb["vertical"][v]["followup"]] for v in vk],
    ])

    # ── SLIDE 3: VERTICALS PIE PAIR ───────────────────────────────────────────
    s3 = CM.add_slide(prs)
    CM.add_logo(s3)
    v_labels = [C.VERTICALS[v] for v in vk]
    va_values = [va["vertical"][v]["total"] for v in vk]
    vb_values = [vb["vertical"][v]["total"] for v in vk]
    CM.add_pie_pair(s3, v_labels, v_labels, va_values, vb_values,
                    ["VERTICALS", la], ["VERTICALS", lb])

    # ── SLIDE 4: NEW vs FOLLOW-UP BAR CHARTS (SIDE BY SIDE) ──────────────────
    s4 = CM.add_slide(prs)
    CM.add_logo(s4)
    CM.add_section_header(s4, "NEW vs FOLLOW-UP BY VERTICAL", Inches(0.15), Inches(0.75),
                          Inches(13.0), Inches(0.45), bg_color=C.PRIMARY)
    CM.add_text(s4, la, Inches(0.15), Inches(1.25), Inches(6.2), Inches(0.30),
                font_size=C.HEADING_FONT, bold=True, color=C.PRIMARY, alignment=PP_ALIGN.CENTER)
    CM.add_text(s4, lb, Inches(6.95), Inches(1.25), Inches(6.2), Inches(0.30),
                font_size=C.HEADING_FONT, bold=True, color=C.PRIMARY, alignment=PP_ALIGN.CENTER)
    CM.add_column_chart(s4, Inches(0.05), Inches(1.60), Inches(6.2), Inches(5.5),
                        v_labels,
                        [("New", [va["vertical"][v]["new"] for v in vk]),
                         ("Follow-up", [va["vertical"][v]["followup"] for v in vk])],
                        legend_position="bottom",
                        series_colors=[C.NEW_CASE_CLR, C.FOLLOWUP_CLR])
    CM.add_column_chart(s4, Inches(6.85), Inches(1.60), Inches(6.2), Inches(5.5),
                        v_labels,
                        [("New", [vb["vertical"][v]["new"] for v in vk]),
                         ("Follow-up", [vb["vertical"][v]["followup"] for v in vk])],
                        legend_position="bottom",
                        series_colors=[C.NEW_CASE_CLR, C.FOLLOWUP_CLR])

    # ── SLIDE 5: STAKEHOLDER PIE PAIR ─────────────────────────────────────────
    s5 = CM.add_slide(prs)
    CM.add_logo(s5)
    st_labels = C.STAKEHOLDER_LABELS
    st_a = [va["stakeholder"].get(s, 0) for s in st_labels]
    st_b = [vb["stakeholder"].get(s, 0) for s in st_labels]
    CM.add_pie_pair(s5, st_labels, st_labels, st_a, st_b,
                    ["STAKEHOLDER", la], ["STAKEHOLDER", lb])

    # ── SLIDE 6: CONCERN PIE PAIR ─────────────────────────────────────────────
    s6 = CM.add_slide(prs)
    CM.add_logo(s6)
    c_labels = C.CONCERN_LABELS
    c_a = [va["concern"].get(c, 0) for c in c_labels]
    c_b = [vb["concern"].get(c, 0) for c in c_labels]
    CM.add_pie_pair(s6, c_labels, c_labels, c_a, c_b,
                    ["RANGE OF CONCERN", la], ["RANGE OF CONCERN", lb])

    # ── SLIDE 7: GENDER & MODE PIE PAIR ───────────────────────────────────────
    s7 = CM.add_slide(prs)
    CM.add_logo(s7)
    g_labels = C.GENDER_LABELS
    g_a = [va["gender"].get(g, 0) for g in g_labels]
    g_b = [vb["gender"].get(g, 0) for g in g_labels]
    m_labels = C.MODE_LABELS
    m_a = [va["mode"].get(m, 0) for m in m_labels]
    m_b = [vb["mode"].get(m, 0) for m in m_labels]
    CM.add_pie_pair(s7, g_labels, m_labels, g_a, m_b,
                    ["GENDER DISTRIBUTION", la], ["MODE OF SESSION", lb])

    # ── SLIDE 8: COMPARISON BAR CHARTS ────────────────────────────────────────
    s8 = CM.add_slide(prs)
    CM.add_logo(s8)
    CM.add_full_column(s8, f"Concern Comparison  |  {la} vs {lb}",
                       c_labels, [(la, c_a), (lb, c_b)],
                       legend_position="bottom")

    s9 = CM.add_slide(prs)
    CM.add_logo(s9)
    CM.add_full_column(s9, f"Stakeholder Comparison  |  {la} vs {lb}",
                       st_labels, [(la, st_a), (lb, st_b)],
                       legend_position="bottom")

    # ── SLIDES 10-12: TEMPLATE (GENDER / CONCERN / STAKEHOLDER, vertical-wise) ─
    TS.add_template_slides(prs, va, vb, label_a=la, label_b=lb)

    CM.add_proposed_points(prs, proposed_points)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


from pptx.enum.text import PP_ALIGN
